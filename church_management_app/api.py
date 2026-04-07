from rest_framework import viewsets
from rest_framework.decorators import action
from rest_framework.permissions import AllowAny, IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.exceptions import PermissionDenied
from rest_framework.parsers import FormParser, JSONParser, MultiPartParser

from django.conf import settings
from django.http import HttpResponse
from django.http import FileResponse
from django.db import transaction
from django.core.files.base import ContentFile
from django.core.mail import EmailMessage
from django.utils import timezone

import datetime
import hashlib

import os
import io
import csv
import re
import secrets

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None
    Inches = None
    Pt = None

import qrcode
import qrcode.image.svg

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from django.db.models import Q, Sum, Count
from django.db.models.functions import TruncDay, TruncMonth, TruncWeek, TruncYear

from .models import (
    Announcement,
    AnnouncementComment,
    AnnouncementCommentLike,
    AnnouncementDeck,
    AnnouncementDeckItem,
    AnnouncementLike,
    ApprovalRequest,
    AuditLogEntry,
    Attendance,
    ChurchBiography,
    ChurchConsistory,
    Contact,
    Department,
    Document,
    Event,
    EventComment,
    EventAttendanceAggregate,
    EventVisitorAggregate,
    EventLogisticsConsumption,
    BaptismEvent,
    BaptismCandidate,
    EvangelismActivity,
    TrainingEvent,
    MarriageRecord,
    Family,
    FinancialCategory,
    FinancialDocumentSequence,
    FinancialTransaction,
    HomeGroup,
    LogisticsItem,
    Member,
    Ministry,
    ActivityDuration,
    Notification,
    ReportCertificate,
    User,
)
from .permissions import (
    IsAdminOrSuperAdmin,
    IsAdminOrSuperAdminOrReadOnly,
    PublicReadAdminWrite,
    IsDepartmentHeadOrAdmin,
    IsEvangelismHeadOrAdmin,
    IsLogisticsHeadOrAdmin,
    IsSecretaryOrAdmin,
    IsTreasurerOrAdmin,
)
from .serializers import (
    AnnouncementSerializer,
    AnnouncementCommentSerializer,
    AnnouncementDeckSerializer,
    AnnouncementDeckItemSerializer,
    ApprovalRequestSerializer,
    AuditLogEntrySerializer,
    AttendanceSerializer,
    ChurchBiographySerializer,
    ChurchConsistorySerializer,
    ContactSerializer,
    DepartmentSerializer,
    DocumentSerializer,
    ActivityDurationSerializer,
    EventAttendanceAggregateSerializer,
    EventVisitorAggregateSerializer,
    EventLogisticsConsumptionSerializer,
    BaptismEventSerializer,
    BaptismCandidateSerializer,
    EvangelismActivitySerializer,
    TrainingEventSerializer,
    MarriageRecordSerializer,
    EventCommentSerializer,
    EventSerializer,
    FamilySerializer,
    FinancialCategorySerializer,
    FinancialTransactionSerializer,
    HomeGroupSerializer,
    LogisticsItemSerializer,
    MeUpdateSerializer,
    MemberSerializer,
    MinistrySerializer,
    NotificationSerializer,
    UserSerializer,
)


def _is_admin_user(user):
    if not user or not getattr(user, 'is_authenticated', False):
        return False
    if getattr(user, 'is_superuser', False):
        return True
    if getattr(user, 'is_staff', False):
        return True
    return getattr(user, 'role', None) in {'super_admin', 'pastor', 'admin', 'administrator'}


def _admin_recipients_qs():
    return User.objects.filter(
        Q(is_superuser=True)
        | Q(is_staff=True)
        | Q(role__in=['super_admin', 'pastor', 'admin', 'administrator'])
    ).distinct()


def _notify_admins(title, message):
    for u in _admin_recipients_qs():
        Notification.objects.create(title=title, message=message, recipient=u)


def _notify_user(user, title, message):
    if not user or not getattr(user, 'is_authenticated', False):
        return
    Notification.objects.create(title=title, message=message, recipient=user)


def _create_approval_request(request, model, action, payload=None, target_object_id=None, object_repr=None):
    ar = ApprovalRequest.objects.create(
        model=str(model or '')[:100],
        action=str(action or '')[:10],
        target_object_id=str(target_object_id)[:64] if target_object_id else None,
        object_repr=str(object_repr)[:200] if object_repr else None,
        payload=payload,
        requested_by=request.user,
        status='pending',
    )

    who = getattr(request.user, 'username', None) or 'Utilisateur'
    label = f"{model} ({action})".strip()
    msg = f"{who} a soumis une action à approuver: {label}.".strip()
    _notify_admins('Action à approuver', msg)
    return ar


def _apply_approval_request(ar, request):
    payload = getattr(ar, 'payload', None) or {}
    model = (getattr(ar, 'model', None) or '').strip()
    action = (getattr(ar, 'action', None) or '').strip()

    if model == 'FinancialTransaction':
        v = FinancialTransactionViewSet()
        v.request = request

        if action == 'create':
            ser = FinancialTransactionSerializer(data=payload)
            ser.is_valid(raise_exception=True)
            with transaction.atomic():
                tx = ser.save(cashier=request.user, created_by=getattr(ar, 'requested_by', None))
                v._ensure_document_number(tx)
                if tx.direction == 'in':
                    v._ensure_receipt_pdf(tx, request)
                    if tx.donor_email:
                        v._email_receipt(tx, request)
                AuditLogEntry.objects.create(
                    actor=request.user,
                    action='create',
                    model='FinancialTransaction',
                    object_id=str(tx.pk),
                    object_repr=tx.document_number or tx.receipt_code or str(tx.pk),
                    ip_address=_client_ip(request),
                    payload=payload,
                )
            return

        if action == 'update':
            obj_id = getattr(ar, 'target_object_id', None)
            tx = FinancialTransaction.objects.filter(id=obj_id).first()
            if not tx:
                raise PermissionDenied('Objet introuvable.')
            ser = FinancialTransactionSerializer(tx, data=payload, partial=True)
            ser.is_valid(raise_exception=True)
            ser.save()
            AuditLogEntry.objects.create(
                actor=request.user,
                action='update',
                model='FinancialTransaction',
                object_id=str(tx.pk),
                object_repr=tx.document_number or tx.receipt_code or str(tx.pk),
                ip_address=_client_ip(request),
                payload=payload,
            )
            return

        if action == 'delete':
            obj_id = getattr(ar, 'target_object_id', None)
            tx = FinancialTransaction.objects.filter(id=obj_id).first()
            if not tx:
                raise PermissionDenied('Objet introuvable.')
            object_id = str(tx.pk)
            object_repr = tx.document_number or tx.receipt_code or str(tx.pk)
            tx.delete()
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model='FinancialTransaction',
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
            return

    if model == 'LogisticsItem':
        if action == 'create':
            ser = LogisticsItemSerializer(data=payload)
            ser.is_valid(raise_exception=True)
            obj = ser.save()
            AuditLogEntry.objects.create(
                actor=request.user,
                action='create',
                model='LogisticsItem',
                object_id=str(obj.pk),
                object_repr=getattr(obj, 'name', None) or str(obj.pk),
                ip_address=_client_ip(request),
                payload=payload,
            )
            return

        if action == 'update':
            obj_id = getattr(ar, 'target_object_id', None)
            obj = LogisticsItem.objects.filter(id=obj_id).first()
            if not obj:
                raise PermissionDenied('Objet introuvable.')
            ser = LogisticsItemSerializer(obj, data=payload, partial=True)
            ser.is_valid(raise_exception=True)
            obj = ser.save()
            AuditLogEntry.objects.create(
                actor=request.user,
                action='update',
                model='LogisticsItem',
                object_id=str(obj.pk),
                object_repr=getattr(obj, 'name', None) or str(obj.pk),
                ip_address=_client_ip(request),
                payload=payload,
            )
            return

        if action == 'delete':
            obj_id = getattr(ar, 'target_object_id', None)
            obj = LogisticsItem.objects.filter(id=obj_id).first()
            if not obj:
                raise PermissionDenied('Objet introuvable.')
            object_id = str(obj.pk)
            object_repr = getattr(obj, 'name', None) or str(obj.pk)
            if getattr(obj, 'is_active', True):
                obj.is_active = False
                obj.save(update_fields=['is_active', 'updated_at'])
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model='LogisticsItem',
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
            return

    if model in {'EvangelismActivity', 'TrainingEvent'}:
        if model == 'EvangelismActivity':
            Obj = EvangelismActivity
            Ser = EvangelismActivitySerializer
            event_type = 'evangelism'
            default_title = 'Évangélisation'
        else:
            Obj = TrainingEvent
            Ser = TrainingEventSerializer
            event_type = 'training'
            default_title = 'Affermissement'

        if action == 'create':
            ser = Ser(data=payload)
            ser.is_valid(raise_exception=True)

            title = (payload.get('title') or default_title)
            date = payload.get('date')
            time = payload.get('time')
            location = payload.get('location')
            moderator = payload.get('moderator')
            trainer = payload.get('trainer')
            lesson = payload.get('lesson')

            ev = Event.objects.create(
                title=str(title)[:200],
                date=date,
                time=time,
                location=location,
                moderator=(str(moderator).strip() if moderator is not None else None) or None,
                event_type=event_type,
                duration_type='daily',
                is_published=True,
                published_at=timezone.now(),
            )
            _ensure_event_share_slug(ev)

            save_kwargs = {'created_by': getattr(ar, 'requested_by', None), 'published_event': ev}
            if model == 'TrainingEvent':
                if trainer is not None:
                    save_kwargs['trainer'] = trainer
                if lesson is not None:
                    save_kwargs['lesson'] = lesson
            obj = ser.save(**save_kwargs)

            AuditLogEntry.objects.create(
                actor=request.user,
                action='create',
                model=model,
                object_id=str(obj.pk),
                object_repr=getattr(obj, 'title', None) or str(obj.pk),
                ip_address=_client_ip(request),
                payload=payload,
            )
            return

        if action == 'update':
            obj_id = getattr(ar, 'target_object_id', None)
            obj = Obj.objects.filter(id=obj_id).first()
            if not obj:
                raise PermissionDenied('Objet introuvable.')
            ser = Ser(obj, data=payload, partial=True)
            ser.is_valid(raise_exception=True)
            obj = ser.save()

            ev = getattr(obj, 'published_event', None)
            if ev is None:
                ev = Event.objects.create(
                    title=(getattr(obj, 'title', None) or default_title)[:200],
                    date=getattr(obj, 'date', None),
                    time=getattr(obj, 'time', None),
                    location=getattr(obj, 'location', None),
                    moderator=getattr(obj, 'moderator', None),
                    event_type=event_type,
                    duration_type='daily',
                    is_published=True,
                    published_at=timezone.now(),
                )
                _ensure_event_share_slug(ev)
                obj.published_event = ev
                obj.save(update_fields=['published_event', 'updated_at'])
            else:
                ev.title = (getattr(obj, 'title', None) or ev.title or default_title)[:200]
                ev.date = getattr(obj, 'date', None)
                ev.time = getattr(obj, 'time', None)
                ev.location = getattr(obj, 'location', None)
                ev.moderator = getattr(obj, 'moderator', None) or ev.moderator
                ev.event_type = event_type
                ev.duration_type = 'daily'
                ev.is_published = True
                if not ev.published_at:
                    ev.published_at = timezone.now()
                ev.save()
                _ensure_event_share_slug(ev)

            AuditLogEntry.objects.create(
                actor=request.user,
                action='update',
                model=model,
                object_id=str(obj.pk),
                object_repr=getattr(obj, 'title', None) or str(obj.pk),
                ip_address=_client_ip(request),
                payload=payload,
            )
            return

        if action == 'delete':
            obj_id = getattr(ar, 'target_object_id', None)
            obj = Obj.objects.filter(id=obj_id).first()
            if not obj:
                raise PermissionDenied('Objet introuvable.')
            object_id = str(obj.pk)
            object_repr = getattr(obj, 'title', None) or str(obj.pk)
            obj.delete()
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model=model,
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
            return

    if model == 'EventAttendanceAggregate' and action == 'update':
        event_id = getattr(ar, 'target_object_id', None)
        event = Event.objects.filter(id=event_id).first()
        if not event:
            raise PermissionDenied('Événement introuvable.')
        agg, _ = EventAttendanceAggregate.objects.get_or_create(event=event)
        ser = EventAttendanceAggregateSerializer(agg, data=payload, partial=True)
        ser.is_valid(raise_exception=True)
        ser.save(updated_by=getattr(ar, 'requested_by', None))
        return

    if model == 'EventVisitorAggregate' and action == 'update':
        event_id = getattr(ar, 'target_object_id', None)
        event = Event.objects.filter(id=event_id).first()
        if not event:
            raise PermissionDenied('Événement introuvable.')
        agg, _ = EventVisitorAggregate.objects.get_or_create(event=event)
        ser = EventVisitorAggregateSerializer(agg, data=payload, partial=True)
        ser.is_valid(raise_exception=True)
        ser.save(updated_by=getattr(ar, 'requested_by', None))
        return

    if model == 'DepartmentCheckin' and action == 'update':
        event_id = getattr(ar, 'target_object_id', None)
        event = Event.objects.filter(id=event_id).first()
        if not event:
            raise PermissionDenied('Événement introuvable.')

        member_id = payload.get('member_id')
        if not member_id:
            raise PermissionDenied('member_id requis.')

        attended = payload.get('attended', True)
        if isinstance(attended, str):
            attended = attended.strip().lower() not in {'0', 'false', 'no'}
        attended = bool(attended)

        dept = getattr(event, 'department', None)
        member = Member.objects.select_related('user').filter(id=member_id, department=dept).first()
        if not member:
            raise PermissionDenied('Membre introuvable ou non appartenant au département.')

        Attendance.objects.update_or_create(
            event=event,
            member=member,
            defaults={'attended': attended, 'checked_in_at': timezone.now() if attended else None},
        )
        return

    if model == 'EventLogisticsConsumption' and action == 'update':
        event_id = getattr(ar, 'target_object_id', None)
        event = Event.objects.filter(id=event_id).first()
        if not event:
            raise PermissionDenied('Événement introuvable.')

        payload_items = payload.get('items')
        if not isinstance(payload_items, list):
            raise PermissionDenied('items doit être une liste.')

        with transaction.atomic():
            for it in payload_items:
                item_id = it.get('item') or it.get('item_id')
                qty = it.get('quantity_used')
                if not item_id:
                    continue
                try:
                    qty_int = int(qty or 0)
                except Exception:
                    raise PermissionDenied('quantity_used invalide.')
                if qty_int < 0:
                    raise PermissionDenied('quantity_used invalide.')

                li = LogisticsItem.objects.filter(id=item_id).first()
                if not li:
                    raise PermissionDenied(f"Matériel introuvable: {item_id}")
                if qty_int > int(getattr(li, 'quantity', 0) or 0):
                    raise PermissionDenied(f"Quantité supérieure au stock pour: {getattr(li, 'name', 'matériel')}")

                EventLogisticsConsumption.objects.update_or_create(
                    event=event,
                    item=li,
                    defaults={'quantity_used': qty_int, 'updated_by': getattr(ar, 'requested_by', None)},
                )
        return

    raise PermissionDenied('Type de demande non pris en charge.')


def _client_ip(request):
    try:
        xff = request.META.get('HTTP_X_FORWARDED_FOR')
        if xff:
            return xff.split(',')[0].strip()
    except Exception:
        return None
    return request.META.get('REMOTE_ADDR')


def _safe_payload(data):
    if data is None:
        return None
    try:
        if hasattr(data, 'dict'):
            data = data.dict()
    except Exception:
        pass

    if isinstance(data, dict):
        out = {}
        for k, v in data.items():
            try:
                if hasattr(v, 'name'):
                    out[k] = v.name
                elif isinstance(v, (list, tuple)):
                    out[k] = [getattr(x, 'name', str(x)) for x in v]
                else:
                    out[k] = v
            except Exception:
                out[k] = str(v)
        return out

    try:
        return str(data)
    except Exception:
        return None


def _pptx_add_title(prs, title, subtitle=None):
    if Inches is None:
        raise RuntimeError('python-pptx indisponible')
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = str(title or '')
    if subtitle is not None:
        try:
            slide.placeholders[1].text = str(subtitle or '')
        except Exception:
            pass
    return slide


def _pptx_add_announcement_slide(prs, header, text, number=None):
    if Inches is None:
        raise RuntimeError('python-pptx indisponible')
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    if not title_shape:
        return slide

    prefix = f"Annonce {number} — " if number is not None else ''
    title_shape.text = f"{prefix}{header}".strip()

    left = Inches(0.8)
    top = Inches(1.7)
    width = Inches(12.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text or '').strip()
    run.font.size = Pt(34)
    return slide


def _unique_report_code():
    for _ in range(10):
        code = secrets.token_urlsafe(18).replace('-', '').replace('_', '')[:40]
        if not ReportCertificate.objects.filter(code=code).exists():
            return code
    return secrets.token_urlsafe(24).replace('-', '').replace('_', '')[:40]


def _create_report_certificate(report_type, payload, user):
    code = _unique_report_code()
    return ReportCertificate.objects.create(
        code=code,
        report_type=report_type,
        payload=payload,
        pdf_sha256=('0' * 64),
        created_by=user if (user and getattr(user, 'is_authenticated', False)) else None,
    )


def _draw_authenticity_qr(c, verify_url, code, x_left, width, margin, qr_size_mm=26, qr_y_mm=22, text_top_mm=34):
    try:
        qr_img = qrcode.make(verify_url)
        qr_buf = io.BytesIO()
        qr_img.save(qr_buf, format='PNG')
        qr_buf.seek(0)

        qr_size = qr_size_mm * mm
        qr_x = width - margin - qr_size - 10 * mm
        qr_y = qr_y_mm * mm
        c.drawImage(ImageReader(qr_buf), qr_x, qr_y, width=qr_size, height=qr_size, mask='auto')
    except Exception:
        pass

    c.setFillColorRGB(0.07, 0.10, 0.16)
    c.setFont('Helvetica-Bold', 10)
    c.drawString(x_left, text_top_mm * mm, 'Authenticité')
    c.setFillColorRGB(0.20, 0.24, 0.30)
    c.setFont('Helvetica', 9)
    c.drawString(x_left, (text_top_mm - 5) * mm, str(verify_url)[:110])
    c.setFillColorRGB(0.45, 0.50, 0.60)
    c.setFont('Helvetica', 8)
    c.drawString(x_left, (text_top_mm - 10) * mm, f"Code: {str(code)}")


def _ensure_event_share_slug(event):
    if getattr(event, 'share_slug', None):
        return event.share_slug

    for _ in range(6):
        slug = secrets.token_urlsafe(18)
        if not Event.objects.filter(share_slug=slug).exists():
            event.share_slug = slug
            event.save(update_fields=['share_slug'])
            return slug

    slug = secrets.token_urlsafe(24)
    event.share_slug = slug
    event.save(update_fields=['share_slug'])
    return slug


class MeView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get(self, request):
        return Response(UserSerializer(request.user).data)

    def patch(self, request):
        serializer = MeUpdateSerializer(data=request.data, context={'request': request})
        serializer.is_valid(raise_exception=True)
        serializer.save()
        return Response(UserSerializer(request.user).data)


class DashboardSummaryView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        user = getattr(request, 'user', None)
        raw_role = getattr(user, 'role', None)
        role = raw_role
        if raw_role == 'administrator':
            role = 'admin'
        elif raw_role == 'pastor':
            role = 'super_admin'

        is_admin = _is_admin_user(user) or role in {'admin', 'super_admin'}
        can_members = bool(is_admin)
        can_finance = bool(is_admin or role in {'treasurer', 'financial_head'})
        can_attendance = bool(is_admin or role in {'secretary', 'protocol_head'})

        today = timezone.localdate()
        week_start = today - datetime.timedelta(days=6)
        month_start = today.replace(day=1)

        members_total = None
        members_new_week = None
        if can_members:
            members_total = Member.objects.count()
            members_new_week = Member.objects.filter(created_at__date__gte=week_start).count()

        events_week = Event.objects.filter(date__range=[week_start, today]).count()
        events_upcoming = Event.objects.filter(date__gte=today).count()

        alerts_qs = Event.objects.filter(is_alert=True).order_by('date', 'time')[:6]
        alerts = []
        for ev in alerts_qs:
            alerts.append({
                'id': ev.id,
                'title': ev.title,
                'date': ev.date.isoformat() if getattr(ev, 'date', None) else None,
                'time': str(getattr(ev, 'time', '') or '')[:5] or None,
                'message': getattr(ev, 'alert_message', None),
                'is_published': bool(getattr(ev, 'is_published', False)),
            })

        finance_month_totals = {}
        if can_finance:
            finance_month = FinancialTransaction.objects.filter(date__gte=month_start)
            finance_month_rows = list(
                finance_month.values('currency', 'direction')
                .annotate(total=Sum('amount'))
                .order_by('currency', 'direction')
            )
            for r in finance_month_rows:
                cur = r.get('currency') or 'CDF'
                finance_month_totals.setdefault(cur, {'in': 0.0, 'out': 0.0, 'net': 0.0})
                direction = r.get('direction')
                finance_month_totals[cur][direction] = float(r.get('total') or 0)
            for cur, agg in finance_month_totals.items():
                agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)

        attendance_total = None
        attendance_present = None
        attendance_rate = None
        if can_attendance:
            attendance_week = Attendance.objects.filter(event__date__range=[week_start, today])
            attendance_total = attendance_week.count()
            attendance_present = attendance_week.filter(attended=True).count()
            attendance_rate = (float(attendance_present) / float(attendance_total)) if attendance_total else 0.0

        activities = []
        if can_members:
            for m in Member.objects.select_related('user').order_by('-created_at')[:6]:
                u = getattr(m, 'user', None)
                name = (u.get_full_name() if u else '').strip() or (u.username if u else '') or (m.member_number or '')
                activities.append({
                    'kind': 'member',
                    'title': 'Nouveau membre inscrit',
                    'description': name,
                    'created_at': m.created_at.isoformat() if getattr(m, 'created_at', None) else None,
                })

        if can_finance:
            for tx in FinancialTransaction.objects.select_related('member', 'member__user').order_by('-created_at')[:6]:
                who = tx.donor_name
                if not who and tx.member and tx.member.user:
                    who = tx.member.user.get_full_name().strip() or tx.member.user.username
                who = who or tx.recipient_name or '—'
                label = 'Entrée' if tx.direction == 'in' else 'Sortie'
                activities.append({
                    'kind': 'finance',
                    'title': f"{label} enregistrée",
                    'description': f"{who} • {tx.amount} {tx.currency}",
                    'created_at': tx.created_at.isoformat() if getattr(tx, 'created_at', None) else None,
                })

        for ev in Event.objects.order_by('-created_at')[:6]:
            activities.append({
                'kind': 'event',
                'title': 'Événement',
                'description': f"{ev.title} • {ev.date.isoformat() if getattr(ev, 'date', None) else ''}",
                'created_at': ev.created_at.isoformat() if getattr(ev, 'created_at', None) else None,
            })

        def sort_key(a):
            try:
                return a.get('created_at') or ''
            except Exception:
                return ''

        activities.sort(key=sort_key, reverse=True)
        activities = [a for a in activities if a.get('created_at')][:12]

        return Response({
            'permissions': {
                'members': can_members,
                'finance': can_finance,
                'attendance': can_attendance,
            },
            'stats': {
                'members_total': members_total,
                'members_new_week': members_new_week,
                'events_week': events_week,
                'events_upcoming': events_upcoming,
                'alerts_count': Event.objects.filter(is_alert=True).count(),
                'finance_month_totals': finance_month_totals,
                'attendance_rate_week': attendance_rate,
                'attendance_present_week': attendance_present,
                'attendance_rows_week': attendance_total,
                # Nouveaux stats pour Contact, ChurchBiography et ChurchConsistory
                'contacts_new': Contact.objects.filter(status='new').count(),
                'contacts_total': Contact.objects.count(),
                'contacts_week': Contact.objects.filter(created_at__date__gte=week_start).count(),
                'church_biography_active': ChurchBiography.objects.filter(is_active=True).count(),
                'church_consistory_active': ChurchConsistory.objects.filter(is_active=True).count(),
            },
            'alerts': alerts,
            'recent_activity': activities,
        })


class UserViewSet(viewsets.ModelViewSet):
    queryset = User.objects.all().order_by('-id')
    serializer_class = UserSerializer
    permission_classes = [IsAdminOrSuperAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    @action(detail=True, methods=['post'], url_path='block')
    def block(self, request, pk=None):
        user = self.get_object()
        if user.is_active is False:
            return Response(UserSerializer(user).data)
        user.is_active = False
        user.save(update_fields=['is_active'])
        return Response(UserSerializer(user).data)

    @action(detail=True, methods=['post'], url_path='unblock')
    def unblock(self, request, pk=None):
        user = self.get_object()
        if user.is_active is True:
            return Response(UserSerializer(user).data)
        user.is_active = True
        user.save(update_fields=['is_active'])
        return Response(UserSerializer(user).data)


class FamilyViewSet(viewsets.ModelViewSet):
    queryset = Family.objects.all().order_by('-id')
    serializer_class = FamilySerializer
    permission_classes = [IsAdminOrSuperAdmin]


class HomeGroupViewSet(viewsets.ModelViewSet):
    queryset = HomeGroup.objects.all().order_by('-id')
    serializer_class = HomeGroupSerializer
    permission_classes = [IsAdminOrSuperAdmin]


class DepartmentViewSet(viewsets.ModelViewSet):
    queryset = Department.objects.all().order_by('-id')
    serializer_class = DepartmentSerializer
    permission_classes = [IsAdminOrSuperAdmin]


class MinistryViewSet(viewsets.ModelViewSet):
    queryset = Ministry.objects.all().order_by('-id')
    serializer_class = MinistrySerializer
    permission_classes = [IsAdminOrSuperAdmin]


class ActivityDurationViewSet(viewsets.ModelViewSet):
    queryset = ActivityDuration.objects.all().order_by('sort_order', 'label', 'id')
    serializer_class = ActivityDurationSerializer
    permission_classes = [IsAdminOrSuperAdminOrReadOnly]

    def get_queryset(self):
        qs = super().get_queryset()
        is_active = (self.request.query_params.get('is_active') or '').strip().lower()
        if is_active in {'1', 'true', 'yes'}:
            qs = qs.filter(is_active=True)
        elif is_active in {'0', 'false', 'no'}:
            qs = qs.filter(is_active=False)
        return qs


class MemberViewSet(viewsets.ModelViewSet):
    queryset = Member.objects.select_related('user').all().order_by('-id')
    serializer_class = MemberSerializer
    permission_classes = [IsAdminOrSuperAdmin]

    def get_queryset(self):
        qs = super().get_queryset()
        q = (self.request.query_params.get('q') or '').strip()
        if q:
            qs = qs.filter(
                Q(user__first_name__icontains=q)
                | Q(user__last_name__icontains=q)
                | Q(user__username__icontains=q)
                | Q(user__email__icontains=q)
                | Q(user__phone__icontains=q)
                | Q(member_number__icontains=q)
            )
        return qs

    def perform_create(self, serializer):
        obj = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='Member',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'member_number', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def perform_update(self, serializer):
        obj = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='Member',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'member_number', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def destroy(self, request, *args, **kwargs):
        obj = self.get_object()
        object_id = str(obj.pk)
        object_repr = getattr(obj, 'member_number', None) or str(obj.pk)
        resp = super().destroy(request, *args, **kwargs)
        if resp.status_code in {200, 202, 204}:
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model='Member',
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
        return resp

    @action(detail=False, methods=['get'], url_path='export')
    def export(self, request):
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.worksheet.table import Table, TableStyleInfo
        except ModuleNotFoundError:
            return Response(
                {
                    'detail': "Le module 'openpyxl' n'est pas installé sur le serveur. Installez-le puis relancez le backend.",
                },
                status=500,
            )

        wb = Workbook()
        ws = wb.active
        ws.title = 'Membres'

        headers = [
            'N° Membre',
            'Prénom',
            'Postnom',
            'Nom',
            'Téléphone',
            'Email',
            'Sexe',
            'Date de naissance',
            'Lieu de naissance',
            'Nationalité',
            'État civil',
            'Profession',
            'Fonction publique',
            'Poste à l’église',
            'Niveau d’étude',
            'Province',
            'Ville',
            'Commune',
            'Quartier',
            'Avenue',
            'N° maison',
            'Actif',
            'Cause inactivité',
            'Créé le',
        ]

        ws.append(headers)

        header_fill = PatternFill('solid', fgColor='1F4E79')
        header_font = Font(bold=True, color='FFFFFF')
        header_alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        thin = Side(style='thin', color='D9D9D9')
        border = Border(left=thin, right=thin, top=thin, bottom=thin)

        ws.row_dimensions[1].height = 28
        for col_idx, _ in enumerate(headers, start=1):
            c = ws.cell(row=1, column=col_idx)
            c.fill = header_fill
            c.font = header_font
            c.alignment = header_alignment
            c.border = border

        date_cols = {8}
        created_cols = {24}

        for m in self.get_queryset().select_related('user'):
            u = getattr(m, 'user', None)
            created_at = getattr(m, 'created_at', None)
            if created_at:
                try:
                    created_at = timezone.localtime(created_at).replace(tzinfo=None)
                except Exception:
                    created_at = None

            ws.append([
                m.member_number or '',
                getattr(u, 'first_name', '') if u else '',
                m.post_name or '',
                getattr(u, 'last_name', '') if u else '',
                getattr(u, 'phone', '') if u else '',
                getattr(u, 'email', '') if u else '',
                m.gender or '',
                getattr(m, 'birth_date', None),
                m.place_of_birth or '',
                m.nationality or '',
                m.marital_status or '',
                m.occupation or '',
                getattr(m, 'public_function', None) or '',
                getattr(m, 'church_position', None) or '',
                m.education_level or '',
                m.province or '',
                m.city or '',
                m.commune or '',
                m.quarter or '',
                m.avenue or '',
                m.house_number or '',
                'Oui' if m.is_active else 'Non',
                getattr(m, 'inactive_reason', None) or '',
                created_at,
            ])

        last_row = ws.max_row
        last_col = ws.max_column

        body_alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
        for row in ws.iter_rows(min_row=2, max_row=last_row, min_col=1, max_col=last_col):
            for cell in row:
                cell.alignment = body_alignment
                cell.border = border

        for r in range(2, last_row + 1):
            for c in date_cols:
                ws.cell(row=r, column=c).number_format = 'dd/mm/yyyy'
            for c in created_cols:
                ws.cell(row=r, column=c).number_format = 'dd/mm/yyyy hh:mm'

        ws.freeze_panes = 'A2'

        if last_row >= 2:
            table_ref = f"A1:{get_column_letter(last_col)}{last_row}"
            table = Table(displayName='MembersTable', ref=table_ref)
            style = TableStyleInfo(
                name='TableStyleMedium9',
                showFirstColumn=False,
                showLastColumn=False,
                showRowStripes=True,
                showColumnStripes=False,
            )
            table.tableStyleInfo = style
            ws.add_table(table)

        for col_idx in range(1, last_col + 1):
            max_len = 0
            for cell in ws[get_column_letter(col_idx)]:
                v = cell.value
                if v is None:
                    continue
                if hasattr(v, 'strftime'):
                    s = v.strftime('%d/%m/%Y')
                else:
                    s = str(v)
                if len(s) > max_len:
                    max_len = len(s)
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max(10, max_len + 2), 45)

        ws.sheet_view.showGridLines = False
        ws.print_title_rows = '1:1'
        ws.page_setup.orientation = 'landscape'
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0

        out = io.BytesIO()
        wb.save(out)
        out.seek(0)

        filename = f"members_{timezone.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        resp = HttpResponse(
            out.getvalue(),
            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        )
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp

    @action(detail=True, methods=['get'], url_path='qr')
    def qr(self, request, pk=None):
        member = self.get_object()
        if not member.member_number:
            member.save()

        payload = f"CPD|{member.member_number}"
        img = qrcode.make(payload, image_factory=qrcode.image.svg.SvgPathImage)
        svg = img.to_string().decode('utf-8')
        return HttpResponse(svg, content_type='image/svg+xml')

    @action(detail=True, methods=['get'], url_path='fiche')
    def fiche(self, request, pk=None):
        member = self.get_object()
        u = getattr(member, 'user', None)

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        w, h = A4

        def _s(v):
            return '' if v is None else str(v)

        def _wrap(text, max_chars=95):
            s = _s(text).strip()
            if not s:
                return ['—']
            out = []
            cur = ''
            for part in s.split(' '):
                if not cur:
                    cur = part
                    continue
                if len(cur) + 1 + len(part) <= max_chars:
                    cur = f"{cur} {part}"
                else:
                    out.append(cur)
                    cur = part
            if cur:
                out.append(cur)
            return out

        title = 'FICHE MEMBRE'
        member_no = member.member_number or f"ID {member.id}"
        full_name = ' '.join([
            _s(getattr(u, 'first_name', '')).strip(),
            _s(getattr(member, 'post_name', '')).strip(),
            _s(getattr(u, 'last_name', '')).strip(),
        ]).strip()

        brand = FinancialTransactionViewSet()._pdf_brand_layout(
            c,
            doc_title=title,
            doc_subtitle="Identification — Document officiel",
            badge_text="MEMBRE",
            badge_rgb=(0.12, 0.31, 0.47),
        )

        margin = brand['margin']
        left = margin
        right = w - margin
        x = brand['x']
        y = brand['y_start']

        # Bloc en-tête: Photo + QR + infos principales
        photo_box = 26 * mm
        qr_box = 26 * mm
        header_h = max(photo_box, qr_box)
        header_y_top = y
        header_y_bottom = header_y_top - header_h

        # Photo membre
        try:
            photo_path = getattr(getattr(u, 'photo', None), 'path', None)
            if photo_path and os.path.exists(photo_path):
                c.setFillColorRGB(0.98, 0.98, 0.99)
                c.roundRect(left, header_y_bottom, photo_box, photo_box, 8, fill=1, stroke=0)
                c.drawImage(
                    ImageReader(photo_path),
                    left,
                    header_y_bottom,
                    photo_box,
                    photo_box,
                    preserveAspectRatio=True,
                    mask='auto',
                )
            else:
                c.setFillColorRGB(0.98, 0.98, 0.99)
                c.roundRect(left, header_y_bottom, photo_box, photo_box, 8, fill=1, stroke=0)
                c.setFillColorRGB(0.35, 0.35, 0.35)
                c.setFont('Helvetica-Bold', 10)
                initials = ''.join([_s(getattr(u, 'first_name', '')).strip()[:1], _s(getattr(u, 'last_name', '')).strip()[:1]]).upper()
                c.drawCentredString(left + photo_box / 2, header_y_bottom + photo_box / 2 - 4, initials or '—')
        except Exception:
            pass

        # QR code membre (payload stable)
        try:
            payload = f"CPD|{member_no}"
            qr_img = qrcode.make(payload)
            try:
                qr_img = qr_img.convert('RGB')
            except Exception:
                pass
            qx = right - qr_box
            c.setFillColorRGB(1, 1, 1)
            c.roundRect(qx, header_y_bottom, qr_box, qr_box, 8, fill=1, stroke=0)
            c.drawImage(
                ImageReader(qr_img),
                qx,
                header_y_bottom,
                qr_box,
                qr_box,
                preserveAspectRatio=True,
                mask='auto',
            )
        except Exception:
            pass

        # Identité (texte)
        text_x = left + photo_box + 8 * mm
        c.setFillColorRGB(0.12, 0.31, 0.47)
        c.setFont('Helvetica-Bold', 14)
        c.drawString(text_x, header_y_top - 3 * mm, full_name or '—')
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 9)
        c.drawString(text_x, header_y_top - 9 * mm, f"N° {member_no}")
        c.drawString(text_x, header_y_top - 14 * mm, f"Téléphone: {_s(getattr(u, 'phone', None) or '—')}")
        c.drawString(text_x, header_y_top - 19 * mm, f"Email: {_s(getattr(u, 'email', None) or '—')}")

        c.setFillColorRGB(0, 0, 0)
        y = header_y_bottom - 8 * mm

        def section(label):
            nonlocal y
            y -= 6 * mm
            c.setFont('Helvetica-Bold', 11)
            c.setFillColorRGB(0.12, 0.31, 0.47)
            c.drawString(left, y, label)
            c.setFillColorRGB(0, 0, 0)
            y -= 4 * mm

        def kv(label, value):
            nonlocal y
            if y < 22 * mm:
                c.showPage()
                brand2 = FinancialTransactionViewSet()._pdf_brand_layout(
                    c,
                    doc_title=title,
                    doc_subtitle="Identification — Document officiel",
                    badge_text="MEMBRE",
                    badge_rgb=(0.12, 0.31, 0.47),
                )
                y = brand2['y_start']
            c.setFont('Helvetica-Bold', 9)
            c.drawString(left, y, f"{label} :")
            c.setFont('Helvetica', 9)
            lines = _wrap(value, max_chars=95)
            first_x = left + 40 * mm
            for i, line in enumerate(lines):
                if y < 22 * mm:
                    c.showPage()
                    brand3 = FinancialTransactionViewSet()._pdf_brand_layout(
                        c,
                        doc_title=title,
                        doc_subtitle="Identification — Document officiel",
                        badge_text="MEMBRE",
                        badge_rgb=(0.12, 0.31, 0.47),
                    )
                    y = brand3['y_start']
                    c.setFont('Helvetica', 9)
                c.drawString(first_x, y, line if i > 0 else line)
                if i < len(lines) - 1:
                    y -= 4 * mm
            y -= 5 * mm

        section('IDENTITÉ')
        kv('Nom complet', full_name or '—')
        kv('Sexe', member.gender or '—')
        kv('Date de naissance', member.birth_date.strftime('%d/%m/%Y') if getattr(member, 'birth_date', None) else '—')
        kv('Lieu de naissance', member.place_of_birth or '—')
        kv('Nationalité', member.nationality or '—')
        kv('État civil', member.marital_status or '—')

        section('CONTACT')
        kv('Téléphone', getattr(u, 'phone', None) or '—')
        kv('Email', getattr(u, 'email', None) or '—')

        section('PROFESSION')
        kv('Profession', member.occupation or '—')
        kv('Fonction publique', getattr(member, 'public_function', None) or '—')
        kv('Poste à l’église', getattr(member, 'church_position', None) or '—')
        kv('Niveau d’étude', member.education_level or '—')

        section('ADRESSE')
        kv('Province', member.province or '—')
        kv('Ville', member.city or '—')
        kv('Commune', member.commune or '—')
        kv('Quartier', member.quarter or '—')
        kv('Avenue', member.avenue or '—')
        kv('N° maison', member.house_number or '—')

        section('INFORMATIONS ÉGLISE')
        kv('Famille', _s(getattr(getattr(member, 'family', None), 'name', None)) or '—')
        kv('Groupe', _s(getattr(getattr(member, 'home_group', None), 'name', None)) or '—')
        kv('Département', _s(getattr(getattr(member, 'department', None), 'name', None)) or '—')
        kv('Ministère', _s(getattr(getattr(member, 'ministry', None), 'name', None)) or '—')
        kv('Date de baptême', member.baptism_date.strftime('%d/%m/%Y') if getattr(member, 'baptism_date', None) else '—')

        section('STATUT')
        kv('Actif', 'Oui' if member.is_active else 'Non')
        kv('Cause inactivité', getattr(member, 'inactive_reason', None) or '—')

        section('URGENCE')
        kv('Nom', member.emergency_contact_name or '—')
        kv('Téléphone', member.emergency_contact_phone or '—')
        kv('Lien', member.emergency_contact_relation or '—')

        created_at = getattr(member, 'created_at', None)
        if created_at:
            try:
                created_at = timezone.localtime(created_at).strftime('%d/%m/%Y %H:%M')
            except Exception:
                created_at = None
        c.setFont('Helvetica', 8)
        c.setFillColorRGB(0.35, 0.35, 0.35)
        c.drawRightString(right, 12 * mm, f"Créé le: {created_at or '—'}")

        c.showPage()
        c.save()
        pdf = buf.getvalue()
        buf.close()

        safe_no = re.sub(r'[^A-Za-z0-9_-]+', '_', member_no)
        filename = f"fiche_membre_{safe_no}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        resp = HttpResponse(pdf, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp


class BaptismEventViewSet(viewsets.ModelViewSet):
    queryset = BaptismEvent.objects.select_related('event').all().order_by('-event__date', '-event__time', '-id')
    serializer_class = BaptismEventSerializer
    permission_classes = [IsSecretaryOrAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def perform_create(self, serializer):
        request = getattr(self, 'request', None)
        user = getattr(request, 'user', None) if request else None

        event = None
        try:
            event_id = request.data.get('event_id') if request else None
        except Exception:
            event_id = None

        if event_id:
            event = Event.objects.filter(id=event_id).first()

        if not event:
            title = (request.data.get('title') or 'Baptême').strip() if request else 'Baptême'
            date = request.data.get('date') if request else None
            time = request.data.get('time') if request else None
            location = request.data.get('location') if request else None
            moderator = (request.data.get('moderator') or '').strip() if request else ''

            event = Event.objects.create(
                title=title,
                date=date,
                time=time,
                location=location,
                moderator=moderator or None,
                event_type='baptism',
                duration_type='daily',
                is_published=True,
                published_at=timezone.now(),
            )
            _ensure_event_share_slug(event)
        else:
            try:
                moderator = (request.data.get('moderator') or '').strip() if request else ''
            except Exception:
                moderator = ''
            if moderator:
                event.moderator = moderator
                event.save(update_fields=['moderator'])

        obj = serializer.save(event=event, created_by=user if (user and user.is_authenticated) else None)
        AuditLogEntry.objects.create(
            actor=user if (user and user.is_authenticated) else None,
            action='create',
            model='BaptismEvent',
            object_id=str(obj.pk),
            object_repr=getattr(event, 'title', None) or str(obj.pk),
            ip_address=_client_ip(request) if request else None,
            payload=_safe_payload(getattr(request, 'data', None)) if request else None,
        )

    @action(detail=True, methods=['get'], url_path='report-pdf')
    def report_pdf(self, request, pk=None):
        be = self.get_object()
        event = getattr(be, 'event', None)
        candidates = BaptismCandidate.objects.filter(baptism_event=be).order_by('name', 'post_name', 'id')

        stamp = timezone.now().strftime('%Y%m%d_%H%M%S')
        filename = f"rapport_bapteme_{getattr(event, 'id', pk)}_{stamp}.pdf"

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        meta = FinancialTransactionViewSet()._pdf_brand_layout(
            c,
            doc_title='Rapport',
            doc_subtitle='Document officiel — baptême',
            badge_text='BAPTÊME',
            badge_rgb=(0.18, 0.55, 0.40),
        )

        width = meta['width']
        margin = meta['margin']
        x = meta['x']
        y = meta['y_start']

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 16)
        c.drawString(x, y, (getattr(event, 'title', None) or 'Baptême')[:70])
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica', 10)
        c.drawRightString(width - margin - 10 * mm, y, f"{getattr(event, 'date', '')}  •  {str(getattr(event, 'time', '') or '')[:5] or '—'}")
        y -= 8 * mm

        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 10)
        c.drawString(x, y, f"Lieu: {getattr(event, 'location', None) or '—'}"[:120])
        y -= 10 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, f"Candidats ({candidates.count()})")
        y -= 8 * mm

        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(x, y, 'Nom')
        c.drawString(x + 70 * mm, y, 'Naissance')
        c.drawRightString(width - margin - 10 * mm, y, 'Lieu')
        y -= 6 * mm

        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 9)
        for cand in candidates:
            label = f"{(cand.name or '').strip()} {(cand.post_name or '').strip()}".strip() or '—'
            bdate = cand.birth_date.isoformat() if getattr(cand, 'birth_date', None) else '—'
            place = (cand.place_of_birth or '—')

            c.drawString(x, y, label[:40])
            c.drawString(x + 70 * mm, y, bdate)
            c.drawRightString(width - margin - 10 * mm, y, place[:40])
            y -= 5.8 * mm
            if y < 25 * mm:
                c.showPage()
                meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                    c,
                    doc_title='Rapport',
                    doc_subtitle='Document officiel — baptême',
                    badge_text='BAPTÊME',
                    badge_rgb=(0.18, 0.55, 0.40),
                )
                width = meta2['width']
                margin = meta2['margin']
                x = meta2['x']
                y = meta2['y_start']

        c.showPage()
        c.save()
        pdf = buf.getvalue()
        buf.close()

        resp = HttpResponse(pdf, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp


class BaptismCandidateViewSet(viewsets.ModelViewSet):
    queryset = BaptismCandidate.objects.select_related('baptism_event', 'baptism_event__event').all().order_by('-id')
    serializer_class = BaptismCandidateSerializer
    permission_classes = [IsSecretaryOrAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_queryset(self):
        qs = super().get_queryset()
        baptism_event = (self.request.query_params.get('baptism_event') or '').strip()
        if baptism_event:
            qs = qs.filter(baptism_event_id=baptism_event)
        return qs

    def perform_create(self, serializer):
        obj = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='BaptismCandidate',
            object_id=str(obj.pk),
            object_repr=(f"{obj.name} {obj.post_name}").strip() or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )


class EvangelismActivityViewSet(viewsets.ModelViewSet):
    queryset = EvangelismActivity.objects.all().order_by('-date', '-time', '-id')
    serializer_class = EvangelismActivitySerializer
    permission_classes = [IsEvangelismHeadOrAdmin]

    def create(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().create(request, *args, **kwargs)
        ar = _create_approval_request(request, model='EvangelismActivity', action='create', payload=_safe_payload(request.data))
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def update(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().update(request, *args, **kwargs)
        obj = self.get_object()
        ar = _create_approval_request(
            request,
            model='EvangelismActivity',
            action='update',
            payload=_safe_payload(request.data),
            target_object_id=getattr(obj, 'id', None),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def partial_update(self, request, *args, **kwargs):
        return self.update(request, *args, **kwargs)

    def destroy(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().destroy(request, *args, **kwargs)
        obj = self.get_object()
        ar = _create_approval_request(
            request,
            model='EvangelismActivity',
            action='delete',
            payload=None,
            target_object_id=getattr(obj, 'id', None),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def perform_create(self, serializer):
        request = getattr(self, 'request', None)
        user = getattr(request, 'user', None) if request else None
        title = (request.data.get('title') or 'Évangélisation').strip() if request else 'Évangélisation'
        date = request.data.get('date') if request else None
        time = request.data.get('time') if request else None
        location = request.data.get('location') if request else None
        moderator = (request.data.get('moderator') or '').strip() if request else ''

        ev = Event.objects.create(
            title=title[:200],
            date=date,
            time=time,
            location=location,
            moderator=moderator or None,
            event_type='evangelism',
            duration_type='daily',
            is_published=True,
            published_at=timezone.now(),
        )
        _ensure_event_share_slug(ev)

        obj = serializer.save(created_by=user if (user and user.is_authenticated) else None, published_event=ev)
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='EvangelismActivity',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'title', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def perform_update(self, serializer):
        obj = serializer.save()
        ev = getattr(obj, 'published_event', None)
        if ev is None:
            ev = Event.objects.create(
                title=(getattr(obj, 'title', None) or 'Évangélisation')[:200],
                date=getattr(obj, 'date', None),
                time=getattr(obj, 'time', None),
                location=getattr(obj, 'location', None),
                moderator=getattr(obj, 'moderator', None),
                event_type='evangelism',
                duration_type='daily',
                is_published=True,
                published_at=timezone.now(),
            )
            _ensure_event_share_slug(ev)
            obj.published_event = ev
            obj.save(update_fields=['published_event', 'updated_at'])
        else:
            ev.title = (getattr(obj, 'title', None) or ev.title or 'Évangélisation')[:200]
            ev.date = getattr(obj, 'date', None)
            ev.time = getattr(obj, 'time', None)
            ev.location = getattr(obj, 'location', None)
            ev.moderator = getattr(obj, 'moderator', None) or ev.moderator
            ev.event_type = 'evangelism'
            ev.duration_type = 'daily'
            ev.is_published = True
            if not ev.published_at:
                ev.published_at = timezone.now()
            ev.save()
            _ensure_event_share_slug(ev)

        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='EvangelismActivity',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'title', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )


class TrainingEventViewSet(viewsets.ModelViewSet):
    queryset = TrainingEvent.objects.all().order_by('-date', '-time', '-id')
    serializer_class = TrainingEventSerializer
    permission_classes = [IsEvangelismHeadOrAdmin]

    def create(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().create(request, *args, **kwargs)
        ar = _create_approval_request(request, model='TrainingEvent', action='create', payload=_safe_payload(request.data))
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def update(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().update(request, *args, **kwargs)
        obj = self.get_object()
        ar = _create_approval_request(
            request,
            model='TrainingEvent',
            action='update',
            payload=_safe_payload(request.data),
            target_object_id=getattr(obj, 'id', None),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def partial_update(self, request, *args, **kwargs):
        return self.update(request, *args, **kwargs)

    def destroy(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().destroy(request, *args, **kwargs)
        obj = self.get_object()
        ar = _create_approval_request(
            request,
            model='TrainingEvent',
            action='delete',
            payload=None,
            target_object_id=getattr(obj, 'id', None),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def perform_create(self, serializer):
        request = getattr(self, 'request', None)
        user = getattr(request, 'user', None) if request else None
        title = (request.data.get('title') or 'Affermissement').strip() if request else 'Affermissement'
        date = request.data.get('date') if request else None
        time = request.data.get('time') if request else None
        location = request.data.get('location') if request else None

        ev = Event.objects.create(
            title=title[:200],
            date=date,
            time=time,
            location=location,
            event_type='training',
            duration_type='daily',
            is_published=True,
            published_at=timezone.now(),
        )
        _ensure_event_share_slug(ev)

        obj = serializer.save(created_by=user if (user and user.is_authenticated) else None, published_event=ev)
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='TrainingEvent',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'title', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def perform_update(self, serializer):
        obj = serializer.save()
        ev = getattr(obj, 'published_event', None)
        if ev is None:
            ev = Event.objects.create(
                title=(getattr(obj, 'title', None) or 'Affermissement')[:200],
                date=getattr(obj, 'date', None),
                time=getattr(obj, 'time', None),
                location=getattr(obj, 'location', None),
                event_type='training',
                duration_type='daily',
                is_published=True,
                published_at=timezone.now(),
            )
            _ensure_event_share_slug(ev)
            obj.published_event = ev
            obj.save(update_fields=['published_event', 'updated_at'])
        else:
            ev.title = (getattr(obj, 'title', None) or ev.title or 'Affermissement')[:200]
            ev.date = getattr(obj, 'date', None)
            ev.time = getattr(obj, 'time', None)
            ev.location = getattr(obj, 'location', None)
            ev.event_type = 'training'
            ev.duration_type = 'daily'
            ev.is_published = True
            if not ev.published_at:
                ev.published_at = timezone.now()
            ev.save()
            _ensure_event_share_slug(ev)

        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='TrainingEvent',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'title', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )


class MarriageRecordViewSet(viewsets.ModelViewSet):
    queryset = MarriageRecord.objects.select_related('groom', 'groom__user', 'bride', 'bride__user', 'published_event').all().order_by('-planned_date', '-planned_time', '-id')
    serializer_class = MarriageRecordSerializer
    permission_classes = [IsAuthenticated]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_queryset(self):
        qs = super().get_queryset()
        u = getattr(self.request, 'user', None)
        if not u or not getattr(u, 'is_authenticated', False):
            return qs.none()

        role = getattr(u, 'role', None)
        is_privileged = bool(
            getattr(u, 'is_superuser', False)
            or getattr(u, 'is_staff', False)
            or role in {'super_admin', 'admin', 'secretary'}
        )
        if is_privileged:
            return qs
        return qs.filter(created_by=u)

    def get_permissions(self):
        action = getattr(self, 'action', None)
        if action in {'list', 'retrieve', 'create'}:
            return [IsAuthenticated()]
        return [IsSecretaryOrAdmin()]

    def _maybe_publish_event(self, obj):
        needs_publish = bool(obj.dowry_paid and obj.civil_verified and obj.prenuptial_tests and obj.approved)
        if not needs_publish:
            return

        groom_u = getattr(getattr(obj, 'groom', None), 'user', None)
        bride_u = getattr(getattr(obj, 'bride', None), 'user', None)
        groom_fallback = (getattr(obj, 'groom_full_name', None) or '').strip() or None
        bride_fallback = (getattr(obj, 'bride_full_name', None) or '').strip() or None

        groom_name = (
            (groom_u.get_full_name() if groom_u else '').strip()
            or (groom_u.username if groom_u else None)
            or (getattr(obj.groom, 'member_number', None) if getattr(obj, 'groom', None) else None)
            or groom_fallback
            or '—'
        )
        bride_name = (
            (bride_u.get_full_name() if bride_u else '').strip()
            or (bride_u.username if bride_u else None)
            or (getattr(obj.bride, 'member_number', None) if getattr(obj, 'bride', None) else None)
            or bride_fallback
            or '—'
        )
        title = f"Mariage: {groom_name} & {bride_name}"[:200]

        if obj.published_event:
            ev = obj.published_event
            ev.title = title
            ev.date = obj.planned_date
            ev.time = obj.planned_time
            ev.location = obj.location
            ev.event_type = 'marriage'
            ev.is_published = True
            if not ev.published_at:
                ev.published_at = timezone.now()
            ev.save()
            _ensure_event_share_slug(ev)
            return

        ev = Event.objects.create(
            title=title,
            date=obj.planned_date,
            time=obj.planned_time,
            location=obj.location,
            event_type='marriage',
            duration_type='daily',
            is_published=True,
            published_at=timezone.now(),
        )
        _ensure_event_share_slug(ev)
        obj.published_event = ev
        obj.save(update_fields=['published_event', 'updated_at'])

    def perform_create(self, serializer):
        obj = serializer.save(created_by=self.request.user)
        self._maybe_publish_event(obj)
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='MarriageRecord',
            object_id=str(obj.pk),
            object_repr=str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def perform_update(self, serializer):
        obj = serializer.save()
        self._maybe_publish_event(obj)
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='MarriageRecord',
            object_id=str(obj.pk),
            object_repr=str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )


class EventViewSet(viewsets.ModelViewSet):
    queryset = Event.objects.all().order_by('-date', '-time')
    serializer_class = EventSerializer
    permission_classes = [IsAdminOrSuperAdminOrReadOnly]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_queryset(self):
        qs = super().get_queryset()
        qp = getattr(self.request, 'query_params', None) or {}

        date_str = (qp.get('date') or '').strip()
        time_str = (qp.get('time') or '').strip()
        event_type = (qp.get('event_type') or '').strip()
        department = (qp.get('department') or '').strip()

        if date_str:
            try:
                qs = qs.filter(date=date_str)
            except Exception:
                pass

        if event_type:
            qs = qs.filter(event_type=event_type)

        if department:
            qs = qs.filter(department_id=department)

        if time_str:
            try:
                parts = time_str[:5].split(':')
                h = int(parts[0])
                m = int(parts[1]) if len(parts) > 1 else 0
                target = datetime.time(hour=h, minute=m)
                window_minutes = int((qp.get('window_minutes') or 90))

                base_dt = datetime.datetime.combine(datetime.date(2000, 1, 1), target)
                low = (base_dt - datetime.timedelta(minutes=window_minutes)).time()
                high = (base_dt + datetime.timedelta(minutes=window_minutes)).time()
                if low <= high:
                    qs = qs.filter(time__range=(low, high))
                else:
                    qs = qs.filter(Q(time__gte=low) | Q(time__lte=high))
            except Exception:
                pass

        return qs

    def get_permissions(self):
        if getattr(self, 'action', None) in {'public', 'public_comment'}:
            return [AllowAny()]
        if getattr(self, 'action', None) in {
            'attendance_aggregate',
            'visitor_aggregate',
            'logistics_consumption',
            'department_members',
            'department_checkin',
            'activity_report',
            'activity_report_pdf',
            'attendance_report',
            'validate_closure',
            'set_alert',
            'clear_alert',
        }:
            act = getattr(self, 'action', None)
            if act in {'attendance_aggregate', 'visitor_aggregate', 'department_members', 'department_checkin', 'attendance_report'}:
                return [IsSecretaryOrAdmin()]
            if act in {'logistics_consumption'}:
                return [IsLogisticsHeadOrAdmin()]
            return [IsAdminOrSuperAdmin()]
        return super().get_permissions()

    @action(detail=True, methods=['post'], url_path='validate-closure')
    def validate_closure(self, request, pk=None):
        event = self.get_object()
        event.closure_validated_at = timezone.now()
        event.save(update_fields=['closure_validated_at'])
        return Response(EventSerializer(event, context={'request': request}).data)

    @action(detail=True, methods=['post'], url_path='set-alert')
    def set_alert(self, request, pk=None):
        event = self.get_object()
        message = request.data.get('message')
        event.is_alert = True
        event.alert_message = (str(message).strip() if message is not None else None) or event.alert_message
        event.save(update_fields=['is_alert', 'alert_message'])
        return Response(EventSerializer(event, context={'request': request}).data)

    @action(detail=True, methods=['post'], url_path='clear-alert')
    def clear_alert(self, request, pk=None):
        event = self.get_object()
        event.is_alert = False
        event.alert_message = None
        event.save(update_fields=['is_alert', 'alert_message'])
        return Response(EventSerializer(event, context={'request': request}).data)

    def _ensure_share_slug(self, event):
        if getattr(event, 'share_slug', None):
            return event.share_slug

        for _ in range(6):
            slug = secrets.token_urlsafe(18)
            if not Event.objects.filter(share_slug=slug).exists():
                event.share_slug = slug
                event.save(update_fields=['share_slug'])
                return slug

        slug = secrets.token_urlsafe(24)
        event.share_slug = slug
        event.save(update_fields=['share_slug'])
        return slug

    def perform_create(self, serializer):
        event = serializer.save()
        self._ensure_share_slug(event)
        if getattr(event, 'is_published', False) and not getattr(event, 'published_at', None):
            event.published_at = timezone.now()
            event.save(update_fields=['published_at'])

    def perform_update(self, serializer):
        event = serializer.save()
        self._ensure_share_slug(event)
        if getattr(event, 'is_published', False) and not getattr(event, 'published_at', None):
            event.published_at = timezone.now()
            event.save(update_fields=['published_at'])

    @action(detail=True, methods=['post'], url_path='publish')
    def publish(self, request, pk=None):
        event = self.get_object()
        self._ensure_share_slug(event)
        event.is_published = True
        if not event.published_at:
            event.published_at = timezone.now()
        event.save(update_fields=['is_published', 'published_at', 'share_slug'])
        return Response(EventSerializer(event).data)

    @action(detail=True, methods=['post'], url_path='unpublish')
    def unpublish(self, request, pk=None):
        event = self.get_object()
        event.is_published = False
        event.save(update_fields=['is_published'])
        return Response(EventSerializer(event).data)

    @action(detail=True, methods=['get', 'post'], url_path='attendance-aggregate')
    def attendance_aggregate(self, request, pk=None):
        event = self.get_object()
        agg, _ = EventAttendanceAggregate.objects.get_or_create(event=event)

        if request.method == 'GET':
            return Response(EventAttendanceAggregateSerializer(agg).data)

        if not _is_admin_user(request.user):
            ar = _create_approval_request(
                request,
                model='EventAttendanceAggregate',
                action='update',
                payload=_safe_payload(request.data),
                target_object_id=getattr(event, 'id', None),
                object_repr=getattr(event, 'title', None) or str(getattr(event, 'id', '') or ''),
            )
            return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

        serializer = EventAttendanceAggregateSerializer(agg, data=request.data, partial=True)
        serializer.is_valid(raise_exception=True)
        serializer.save(updated_by=request.user)
        return Response(EventAttendanceAggregateSerializer(agg).data)

    @action(detail=True, methods=['get', 'post'], url_path='visitor-aggregate')
    def visitor_aggregate(self, request, pk=None):
        event = self.get_object()
        agg, _ = EventVisitorAggregate.objects.get_or_create(event=event)

        if request.method == 'GET':
            return Response(EventVisitorAggregateSerializer(agg).data)

        if not _is_admin_user(request.user):
            ar = _create_approval_request(
                request,
                model='EventVisitorAggregate',
                action='update',
                payload=_safe_payload(request.data),
                target_object_id=getattr(event, 'id', None),
                object_repr=getattr(event, 'title', None) or str(getattr(event, 'id', '') or ''),
            )
            return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

        serializer = EventVisitorAggregateSerializer(agg, data=request.data, partial=True)
        serializer.is_valid(raise_exception=True)
        serializer.save(updated_by=request.user)
        return Response(EventVisitorAggregateSerializer(agg).data)

    @action(detail=True, methods=['get', 'post'], url_path='logistics-consumption')
    def logistics_consumption(self, request, pk=None):
        event = self.get_object()

        if request.method == 'GET':
            qs = EventLogisticsConsumption.objects.select_related('item').filter(event=event).order_by('item__name', 'id')
            out = []
            for row in qs:
                item = getattr(row, 'item', None)
                out.append({
                    'id': row.id,
                    'item': row.item_id,
                    'item_name': getattr(item, 'name', None) if item else None,
                    'item_unit': getattr(item, 'unit', None) if item else None,
                    'quantity_used': row.quantity_used,
                })
            return Response({'items': out})

        if not _is_admin_user(request.user):
            ar = _create_approval_request(
                request,
                model='EventLogisticsConsumption',
                action='update',
                payload=_safe_payload(request.data),
                target_object_id=getattr(event, 'id', None),
                object_repr=getattr(event, 'title', None) or str(getattr(event, 'id', '') or ''),
            )
            return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

        payload_items = request.data.get('items')
        if not isinstance(payload_items, list):
            return Response({'detail': 'items doit être une liste.'}, status=400)

        updated = []
        with transaction.atomic():
            for it in payload_items:
                item_id = it.get('item') or it.get('item_id')
                qty = it.get('quantity_used')
                if not item_id:
                    continue
                try:
                    qty_int = int(qty or 0)
                except Exception:
                    return Response({'detail': 'quantity_used invalide.'}, status=400)
                if qty_int < 0:
                    return Response({'detail': 'quantity_used invalide.'}, status=400)

                li = LogisticsItem.objects.filter(id=item_id).first()
                if not li:
                    return Response({'detail': f"Matériel introuvable: {item_id}"}, status=404)
                if qty_int > int(getattr(li, 'quantity', 0) or 0):
                    return Response({'detail': f"Quantité supérieure au stock pour: {getattr(li, 'name', 'matériel')}"}, status=400)

                obj, _ = EventLogisticsConsumption.objects.update_or_create(
                    event=event,
                    item=li,
                    defaults={'quantity_used': qty_int, 'updated_by': request.user},
                )
                updated.append(obj)

        qs = EventLogisticsConsumption.objects.select_related('item').filter(event=event).order_by('item__name', 'id')
        out = []
        for row in qs:
            item = getattr(row, 'item', None)
            out.append({
                'id': row.id,
                'item': row.item_id,
                'item_name': getattr(item, 'name', None) if item else None,
                'item_unit': getattr(item, 'unit', None) if item else None,
                'quantity_used': row.quantity_used,
            })
        return Response({'items': out})

    @action(detail=True, methods=['get'], url_path='department-members')
    def department_members(self, request, pk=None):
        event = self.get_object()
        dept = getattr(event, 'department', None)
        if not dept:
            return Response({'detail': 'Département non défini pour cet événement.'}, status=400)

        qs = Member.objects.select_related('user').filter(department=dept, is_active=True).order_by('user__last_name', 'user__first_name', 'id')
        out = []
        for m in qs:
            u = getattr(m, 'user', None)
            out.append({
                'id': m.id,
                'member_number': getattr(m, 'member_number', None),
                'first_name': getattr(u, 'first_name', None) if u else None,
                'last_name': getattr(u, 'last_name', None) if u else None,
                'username': getattr(u, 'username', None) if u else None,
            })
        return Response({'department_id': dept.id, 'department_name': dept.name, 'members': out})

    @action(detail=True, methods=['post'], url_path='department-checkin')
    def department_checkin(self, request, pk=None):
        event = self.get_object()
        dept = getattr(event, 'department', None)
        if not dept:
            return Response({'detail': 'Département non défini pour cet événement.'}, status=400)

        member_id = request.data.get('member_id')
        if not member_id:
            return Response({'detail': 'member_id requis'}, status=400)

        if not _is_admin_user(request.user):
            ar = _create_approval_request(
                request,
                model='DepartmentCheckin',
                action='update',
                payload=_safe_payload(request.data),
                target_object_id=getattr(event, 'id', None),
                object_repr=getattr(event, 'title', None) or str(getattr(event, 'id', '') or ''),
            )
            return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

        attended = request.data.get('attended', True)
        if isinstance(attended, str):
            attended = attended.strip().lower() not in {'0', 'false', 'no'}
        attended = bool(attended)

        member = Member.objects.select_related('user').filter(id=member_id, department=dept).first()
        if not member:
            return Response({'detail': 'Membre introuvable ou non appartenant au département.'}, status=404)

        att, _ = Attendance.objects.update_or_create(
            event=event,
            member=member,
            defaults={'attended': attended, 'checked_in_at': timezone.now() if attended else None},
        )
        return Response(AttendanceSerializer(att).data)

    def _age_years(self, birth_date):
        if not birth_date:
            return None
        today = timezone.localdate()
        try:
            years = today.year - birth_date.year
            if (today.month, today.day) < (birth_date.month, birth_date.day):
                years -= 1
            return years
        except Exception:
            return None

    def _activity_report_data(self, request, event):
        et = getattr(event, 'event_type', None)
        attendance = None
        visitors = None

        if et not in {'evangelism', 'training'}:
            agg, _ = EventAttendanceAggregate.objects.get_or_create(event=event)
            vis, _ = EventVisitorAggregate.objects.get_or_create(event=event)

            attendance = {
                'anonymous': {
                    'male_adults': agg.male_adults,
                    'female_adults': agg.female_adults,
                    'male_children': agg.male_children,
                    'female_children': agg.female_children,
                    'total': (agg.male_adults + agg.female_adults + agg.male_children + agg.female_children),
                },
                'department': None,
            }

            visitors = {
                'male_visitors': vis.male_visitors,
                'female_visitors': vis.female_visitors,
                'total': (vis.male_visitors + vis.female_visitors),
            }

        cons_qs = EventLogisticsConsumption.objects.select_related('item').filter(event=event).order_by('item__name', 'id')
        consumption = []
        for row in cons_qs:
            item = getattr(row, 'item', None)
            consumption.append({
                'item_id': row.item_id,
                'item_name': getattr(item, 'name', None) if item else None,
                'item_unit': getattr(item, 'unit', None) if item else None,
                'quantity_used': row.quantity_used,
            })

        if attendance is not None and getattr(event, 'department', None):
            child_age = 13
            dept_att = Attendance.objects.select_related('member', 'member__user').filter(event=event, attended=True)
            stats = {
                'members_present': 0,
                'men': 0,
                'women': 0,
                'children': 0,
            }
            for a in dept_att:
                m = getattr(a, 'member', None)
                if not m:
                    continue
                stats['members_present'] += 1
                age = self._age_years(getattr(m, 'birth_date', None))
                if age is not None and age < child_age:
                    stats['children'] += 1
                    continue
                g = getattr(m, 'gender', None)
                if g == 'M':
                    stats['men'] += 1
                elif g == 'F':
                    stats['women'] += 1
            attendance['department'] = {
                'department_id': event.department_id,
                'department_name': getattr(event.department, 'name', None) if getattr(event, 'department', None) else None,
                'stats': stats,
            }

        finance_payload = None
        tx_qs = None
        if et not in {'evangelism', 'training'}:
            tx_qs = FinancialTransaction.objects.filter(event=event)

        if tx_qs is not None:
            rows = (
                tx_qs.values('currency', 'direction')
                .annotate(total=Sum('amount'))
                .order_by('currency', 'direction')
            )
            totals = {}
            for r in rows:
                cur = r['currency'] or 'CDF'
                totals.setdefault(cur, {'in': 0.0, 'out': 0.0, 'net': 0.0})
                direction = r['direction']
                totals[cur][direction] = float(r['total'] or 0)
            for cur, t in totals.items():
                t['net'] = float(t.get('in', 0) or 0) - float(t.get('out', 0) or 0)

            breakdown = []
            by_type = {}
            if et == 'baptism':
                breakdown_rows = (
                    tx_qs.values('currency', 'direction', 'transaction_type', 'description')
                    .annotate(total=Sum('amount'))
                    .order_by('currency', 'direction', 'transaction_type', 'description')
                )
                for r in breakdown_rows:
                    breakdown.append({
                        'currency': r.get('currency') or 'CDF',
                        'direction': r.get('direction') or 'in',
                        'transaction_type': r.get('transaction_type') or None,
                        'description': (r.get('description') or '').strip() or None,
                        'total': float(r.get('total') or 0),
                    })

                for item in breakdown:
                    cur = item.get('currency') or 'CDF'
                    direction = item.get('direction') or 'in'
                    tx_type = item.get('transaction_type') or 'other'
                    by_type.setdefault(cur, {})
                    by_type[cur].setdefault(tx_type, {'in': 0.0, 'out': 0.0, 'net': 0.0})
                    by_type[cur][tx_type][direction] += float(item.get('total') or 0)
                for cur, type_map in by_type.items():
                    for tx_type, agg_map in type_map.items():
                        agg_map['net'] = float(agg_map.get('in', 0) or 0) - float(agg_map.get('out', 0) or 0)

            finance_payload = {
                'totals': totals,
                'transaction_count': tx_qs.count(),
                'breakdown': breakdown,
                'by_type': by_type,
            }

        baptism_payload = None
        if et == 'baptism':
            be = getattr(event, 'baptism_event', None)
            if be:
                cqs = BaptismCandidate.objects.filter(baptism_event=be).order_by('name', 'post_name', 'id')
                baptism_payload = {
                    'executors': getattr(be, 'executors', None) or [],
                    'candidates_count': cqs.count(),
                    'candidates': [
                        {
                            'id': c.id,
                            'name': c.name,
                            'post_name': c.post_name,
                            'birth_date': c.birth_date.isoformat() if getattr(c, 'birth_date', None) else None,
                            'place_of_birth': c.place_of_birth,
                        }
                        for c in cqs[:120]
                    ],
                }

        return {
            'event': EventSerializer(event, context={'request': request}).data,
            'attendance': attendance,
            'visitors': visitors,
            'logistics_consumption': consumption,
            'finance': finance_payload,
            'baptism': baptism_payload,
        }

    @action(detail=True, methods=['get'], url_path='activity-report')
    def activity_report(self, request, pk=None):
        event = self.get_object()
        return Response(self._activity_report_data(request, event))

    def _activity_report_pdf(self, request, event):
        data = self._activity_report_data(request, event)

        cert = None
        verify_url = None
        try:
            evp = {
                'event_id': event.id,
                'date': getattr(event, 'date', None).isoformat() if getattr(event, 'date', None) else None,
                'title': getattr(event, 'title', None),
                'generated_at': timezone.now().isoformat(),
            }
            cert = _create_report_certificate('activity', evp, getattr(request, 'user', None))
            verify_url = request.build_absolute_uri(f"/api/reports/verify/?code={cert.code}")
        except Exception:
            cert = None
            verify_url = None

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        badge = 'ACTIVITÉ'
        badge_rgb = (0.20, 0.35, 0.75)
        try:
            et = getattr(event, 'event_type', None)
            if et == 'baptism':
                badge = 'BAPTÊME'
                badge_rgb = (0.18, 0.55, 0.40)
            elif et == 'evangelism':
                badge = 'ÉVANGÉLISATION'
                badge_rgb = (0.05, 0.65, 0.45)
            elif et == 'training':
                badge = 'AFFERMISSEMENT'
                badge_rgb = (0.34, 0.32, 0.76)
            elif et == 'marriage':
                badge = 'MARIAGE'
                badge_rgb = (0.86, 0.20, 0.35)
        except Exception:
            badge = 'ACTIVITÉ'
            badge_rgb = (0.20, 0.35, 0.75)

        meta = FinancialTransactionViewSet()._pdf_brand_layout(
            c,
            doc_title='Rapport',
            doc_subtitle='Document officiel — activité',
            badge_text=badge,
            badge_rgb=badge_rgb,
        )

        width = meta['width']
        margin = meta['margin']
        x = meta['x']
        y = meta['y_start']

        ev = data.get('event') or {}
        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 16)
        c.drawString(x, y, (ev.get('title') or 'Activité')[:70])
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica', 10)
        right = f"{ev.get('date') or '—'}  •  {str(ev.get('time') or '')[:5] or '—'}"
        c.drawRightString(width - margin - 10 * mm, y, right)
        y -= 10 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 11)
        c.drawString(x, y, 'Programme')
        y -= 7 * mm
        c.setFont('Helvetica', 10)
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.drawString(x, y, f"Type: {ev.get('event_type') or '—'}")
        c.drawRightString(width - margin - 10 * mm, y, f"Lieu: {ev.get('location') or '—'}"[:60])
        y -= 10 * mm

        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(x, y, 'Intervenants')
        y -= 6 * mm
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 10)

        sd = ev.get('special_details') or {}
        interv = []
        if ev.get('event_type') == 'training':
            interv = [
                ('Formateur', sd.get('trainer')),
                ('Leçon', sd.get('lesson')),
            ]
        elif ev.get('event_type') == 'evangelism':
            interv = [
                ('Modérateur', sd.get('moderator') or ev.get('moderator')),
            ]
        elif ev.get('event_type') == 'baptism':
            executors = sd.get('executors') or []
            try:
                executors = [str(x).strip() for x in executors if str(x).strip()]
            except Exception:
                executors = []
            interv = [
                ('Exécutant principal', sd.get('moderator') or ev.get('moderator')),
                ('Exécutants', ', '.join(executors[:5]) if executors else None),
            ]
        else:
            interv = [
                ('Modérateur', ev.get('moderator')),
                ('Prédicateur', ev.get('preacher')),
                ('Chorale', ev.get('choir')),
                ('Protocole', ev.get('protocol_team')),
                ('Technique/Media', ev.get('tech_team')),
                ('Communicateur', ev.get('communicator')),
            ]
        shown = 0
        for label, val in interv:
            v = (val or '').strip()
            if not v:
                continue
            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, f"{label}:")
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            c.drawString(x + 28 * mm, y, str(v)[:95])
            y -= 6 * mm
            shown += 1
            if shown >= 6:
                break
        y -= 4 * mm

        if ev.get('event_type') not in {'evangelism', 'training'}:
            att = (data.get('attendance') or {}).get('anonymous') or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Pointage (anonyme)')
            y -= 7 * mm
            c.setFont('Helvetica', 10)
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.drawString(x, y, f"Hommes adultes: {att.get('male_adults', 0)}")
            c.drawString(x + 55 * mm, y, f"Femmes adultes: {att.get('female_adults', 0)}")
            y -= 6 * mm
            c.drawString(x, y, f"Garçons (enfants): {att.get('male_children', 0)}")
            c.drawString(x + 55 * mm, y, f"Filles (enfants): {att.get('female_children', 0)}")
            c.drawRightString(width - margin - 10 * mm, y, f"Total: {att.get('total', 0)}")
            y -= 10 * mm

            visitors = data.get('visitors') or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Visiteurs')
            y -= 7 * mm
            c.setFont('Helvetica', 10)
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.drawString(x, y, f"Hommes: {visitors.get('male_visitors', 0)}")
            c.drawString(x + 55 * mm, y, f"Femmes: {visitors.get('female_visitors', 0)}")
            c.drawRightString(width - margin - 10 * mm, y, f"Total: {visitors.get('total', 0)}")
            y -= 10 * mm

        if ev.get('event_type') == 'baptism':
            bap = data.get('baptism') or {}
            cands = bap.get('candidates') or []
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, f"Candidats ({bap.get('candidates_count', len(cands)) or 0})")
            y -= 7 * mm
            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, 'Nom')
            c.drawRightString(width - margin - 10 * mm, y, 'Naissance')
            y -= 6 * mm
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 9)
            for cand in cands[:35]:
                label = f"{(cand.get('name') or '').strip()} {(cand.get('post_name') or '').strip()}".strip() or '—'
                bdate = cand.get('birth_date') or '—'
                c.drawString(x, y, label[:55])
                c.drawRightString(width - margin - 10 * mm, y, str(bdate)[:10])
                y -= 5.6 * mm
                if y < 25 * mm:
                    c.showPage()
                    meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                        c,
                        doc_title='Rapport',
                        doc_subtitle='Document officiel — activité',
                        badge_text=badge,
                        badge_rgb=badge_rgb,
                    )
                    width = meta2['width']
                    margin = meta2['margin']
                    x = meta2['x']
                    y = meta2['y_start']
                    c.setFillColorRGB(0.20, 0.24, 0.30)
                    c.setFont('Helvetica', 9)
            y -= 4 * mm

        dept = (data.get('attendance') or {}).get('department')
        if dept:
            stats = (dept.get('stats') or {})
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, f"Pointage (département) — {dept.get('department_name') or '—'}")
            y -= 7 * mm
            c.setFont('Helvetica', 10)
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.drawString(x, y, f"Membres présents: {stats.get('members_present', 0)}")
            c.drawString(x + 55 * mm, y, f"Hommes: {stats.get('men', 0)}")
            c.drawString(x + 90 * mm, y, f"Femmes: {stats.get('women', 0)}")
            c.drawRightString(width - margin - 10 * mm, y, f"Enfants: {stats.get('children', 0)}")
            y -= 10 * mm

        consumption = data.get('logistics_consumption') or []
        if consumption:
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Consommation logistique')
            y -= 7 * mm
            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, 'Matériel')
            c.drawRightString(width - margin - 10 * mm, y, 'Quantité')
            y -= 6 * mm
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)

            for row in consumption:
                label = (row.get('item_name') or '—')
                qty = row.get('quantity_used', 0)
                unit = (row.get('item_unit') or '').strip()
                qty_label = f"{qty} {unit}".strip()
                c.drawString(x, y, str(label)[:80])
                c.drawRightString(width - margin - 10 * mm, y, qty_label[:30])
                y -= 6 * mm
                if y < 25 * mm:
                    c.showPage()
                    meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                        c,
                        doc_title='Rapport',
                        doc_subtitle='Document officiel — activité',
                        badge_text=badge,
                        badge_rgb=badge_rgb,
                    )
                    width = meta2['width']
                    margin = meta2['margin']
                    x = meta2['x']
                    y = meta2['y_start']
                    c.setFillColorRGB(0.20, 0.24, 0.30)
                    c.setFont('Helvetica', 10)

            y -= 4 * mm

        fin_block = data.get('finance') or None
        fin = (fin_block or {}).get('totals') or {}
        by_type = (fin_block or {}).get('by_type') or {}
        breakdown = (fin_block or {}).get('breakdown') or []
        if fin_block:
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Finance (liée à l’activité)')
            y -= 7 * mm
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(x, y, 'Devise')
        c.drawRightString(x + 90 * mm, y, 'Entrées')
        c.drawRightString(x + 120 * mm, y, 'Sorties')
        c.drawRightString(width - margin - 10 * mm, y, 'Solde')
        y -= 6 * mm

        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 10)
        if fin_block:
            for cur, agg in fin.items():
                c.drawString(x, y, str(cur))
                c.drawRightString(x + 90 * mm, y, f"{float(agg.get('in', 0) or 0):.2f}")
                c.drawRightString(x + 120 * mm, y, f"{float(agg.get('out', 0) or 0):.2f}")
                c.drawRightString(width - margin - 10 * mm, y, f"{float(agg.get('net', 0) or 0):.2f}")
                y -= 6 * mm
                if y < 25 * mm:
                    c.showPage()
                    meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                        c,
                        doc_title='Rapport',
                        doc_subtitle='Document officiel — activité',
                        badge_text=badge,
                        badge_rgb=badge_rgb,
                    )
                    width = meta2['width']
                    margin = meta2['margin']
                    x = meta2['x']
                    y = meta2['y_start']
                    c.setFillColorRGB(0.20, 0.24, 0.30)
                    c.setFont('Helvetica', 10)

        if by_type:
            y -= 6 * mm
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Synthèse par catégorie')
            y -= 7 * mm
            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, 'Devise')
            c.drawString(x + 16 * mm, y, 'Catégorie')
            c.drawRightString(x + 100 * mm, y, 'Entrées')
            c.drawRightString(x + 130 * mm, y, 'Sorties')
            c.drawRightString(width - margin - 10 * mm, y, 'Solde')
            y -= 6 * mm

            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 9)
            for cur, type_map in by_type.items():
                for tx_type, agg in (type_map or {}).items():
                    c.drawString(x, y, str(cur))
                    c.drawString(x + 16 * mm, y, str(tx_type)[:26])
                    c.drawRightString(x + 100 * mm, y, f"{float((agg or {}).get('in', 0) or 0):.2f}")
                    c.drawRightString(x + 130 * mm, y, f"{float((agg or {}).get('out', 0) or 0):.2f}")
                    c.drawRightString(width - margin - 10 * mm, y, f"{float((agg or {}).get('net', 0) or 0):.2f}")
                    y -= 5.5 * mm
                    if y < 25 * mm:
                        c.showPage()
                        meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                            c,
                            doc_title='Rapport',
                            doc_subtitle='Document officiel — activité',
                            badge_text='ACTIVITÉ',
                            badge_rgb=(0.20, 0.35, 0.75),
                        )
                        width = meta2['width']
                        margin = meta2['margin']
                        x = meta2['x']
                        y = meta2['y_start']
                        c.setFillColorRGB(0.20, 0.24, 0.30)
                        c.setFont('Helvetica', 9)

        if breakdown:
            y -= 6 * mm
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Détails (catégories)')
            y -= 7 * mm
            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, 'Devise')
            c.drawString(x + 16 * mm, y, 'Catégorie')
            c.drawString(x + 55 * mm, y, 'Description')
            c.drawRightString(width - margin - 10 * mm, y, 'Total')
            y -= 6 * mm

            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 9)
            for item in breakdown:
                cur = str(item.get('currency') or 'CDF')
                tx_type = str(item.get('transaction_type') or '—')
                desc = str(item.get('description') or '—')
                total = float(item.get('total') or 0)
                if total == 0:
                    continue
                c.drawString(x, y, cur)
                c.drawString(x + 16 * mm, y, tx_type[:18])
                c.drawString(x + 55 * mm, y, desc[:36])
                c.drawRightString(width - margin - 10 * mm, y, f"{total:.2f}")
                y -= 5.5 * mm
                if y < 25 * mm:
                    c.showPage()
                    meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                        c,
                        doc_title='Rapport',
                        doc_subtitle='Document officiel — activité',
                        badge_text='ACTIVITÉ',
                        badge_rgb=(0.20, 0.35, 0.75),
                    )
                    width = meta2['width']
                    margin = meta2['margin']
                    x = meta2['x']
                    y = meta2['y_start']
                    c.setFillColorRGB(0.20, 0.24, 0.30)
                    c.setFont('Helvetica', 9)

        if cert and verify_url:
            try:
                _draw_authenticity_qr(c, verify_url, cert.code, x, width, margin, qr_size_mm=26)
            except Exception:
                pass

        c.showPage()
        c.save()
        pdf = buf.getvalue()
        buf.close()

        if cert:
            try:
                cert.pdf_sha256 = hashlib.sha256(pdf).hexdigest()
                cert.save(update_fields=['pdf_sha256', 'updated_at'])
            except Exception:
                pass

        return pdf

    @action(detail=True, methods=['get'], url_path='activity-report-pdf')
    def activity_report_pdf(self, request, pk=None):
        event = self.get_object()
        pdf = self._activity_report_pdf(request, event)
        filename = f"activity_report_{event.id}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        resp = HttpResponse(pdf, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp

    @action(detail=True, methods=['get'], url_path='attendance-report')
    def attendance_report(self, request, pk=None):
        event = self.get_object()
        format_type = request.query_params.get('format', 'pdf').lower()
        
        if format_type not in ['pdf', 'csv', 'xlsx', 'excel']:
            return Response({'error': 'Format must be pdf, csv or xlsx'}, status=400)
        
        if format_type == 'pdf':
            return self._attendance_report_pdf(request, event)
        if format_type in {'xlsx', 'excel'}:
            return self._attendance_report_xlsx(request, event)
        else:
            return self._attendance_report_csv(request, event)

    def _attendance_report_pdf(self, request, event):
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import black, grey
        import logging
        
        logger = logging.getLogger(__name__)
        
        try:
            width, height = A4
            margin = 15 * mm
            
            # Créer le buffer avec une taille suffisante
            buf = io.BytesIO()
            c = canvas.Canvas(buf, pagesize=A4)
            
            # Configuration des polices
            c.setFillColor(black)
            
            # En-tête avec cadre
            c.setStrokeColor(grey)
            c.setLineWidth(1)
            c.rect(margin - 5, height - margin - 70, width - 2*margin + 10, 65, fill=0)
            
            # Titre
            c.setFont("Helvetica-Bold", 16)
            title_text = f"Rapport de Pointage - {event.title or 'Événement'}"
            # Limiter la longueur du titre pour éviter les dépassements
            if len(title_text) > 60:
                title_text = title_text[:57] + "..."
            c.drawString(margin, height - margin - 20, title_text)
            
            # Informations événement
            c.setFont("Helvetica", 11)
            y_pos = height - margin - 40
            
            date_str = event.date.strftime('%d/%m/%Y') if event.date else 'Non spécifiée'
            c.drawString(margin, y_pos, f"Date: {date_str}")
            y_pos -= 15
            
            if event.time:
                time_str = event.time.strftime('%H:%M') if event.time else ''
                c.drawString(margin, y_pos, f"Heure: {time_str}")
                y_pos -= 15
            
            if event.location:
                location_text = event.location[:50] + "..." if len(event.location) > 50 else event.location
                c.drawString(margin, y_pos, f"Lieu: {location_text}")
                y_pos -= 15
            
            # Ligne de séparation
            c.setStrokeColor(grey)
            c.setLineWidth(0.5)
            c.line(margin, y_pos, width - margin, y_pos)
            y_pos -= 25
            
            # Section Participants
            c.setFont("Helvetica-Bold", 14)
            c.drawString(margin, y_pos, "PARTICIPANTS")
            y_pos -= 20
            
            try:
                attendance_agg = EventAttendanceAggregate.objects.filter(event=event).first()
                
                if attendance_agg:
                    c.setFont("Helvetica", 11)
                    
                    # Calculs sécurisés
                    male_adults = max(0, attendance_agg.male_adults or 0)
                    female_adults = max(0, attendance_agg.female_adults or 0)
                    male_children = max(0, attendance_agg.male_children or 0)
                    female_children = max(0, attendance_agg.female_children or 0)
                    
                    total_participants = male_adults + female_adults + male_children + female_children
                    
                    # Tableau des statistiques
                    stats_data = [
                        ("Total participants", total_participants),
                        ("Hommes adultes", male_adults),
                        ("Femmes adultes", female_adults),
                        ("Garçons", male_children),
                        ("Filles", female_children)
                    ]
                    
                    for label, value in stats_data:
                        c.drawString(margin + 10, y_pos, f"{label}: {value}")
                        y_pos -= 15
                else:
                    c.setFont("Helvetica", 11)
                    c.drawString(margin + 10, y_pos, "Aucune donnée de pointage enregistrée")
                    y_pos -= 15
                
            except Exception as agg_error:
                logger.error(f"Erreur lors de la récupération des agrégats: {agg_error}")
                c.setFont("Helvetica", 11)
                c.drawString(margin + 10, y_pos, "Erreur dans les données de participants")
                y_pos -= 15
            
            y_pos -= 10
            
            # Section Visiteurs
            c.setFont("Helvetica-Bold", 14)
            c.drawString(margin, y_pos, "VISITEURS")
            y_pos -= 20
            
            try:
                visitor_agg = EventVisitorAggregate.objects.filter(event=event).first()
                
                if visitor_agg:
                    c.setFont("Helvetica", 11)
                    
                    male_visitors = max(0, visitor_agg.male_visitors or 0)
                    female_visitors = max(0, visitor_agg.female_visitors or 0)
                    total_visitors = male_visitors + female_visitors
                    
                    visitor_stats = [
                        ("Total visiteurs", total_visitors),
                        ("Hommes visiteurs", male_visitors),
                        ("Femmes visiteuses", female_visitors)
                    ]
                    
                    for label, value in visitor_stats:
                        c.drawString(margin + 10, y_pos, f"{label}: {value}")
                        y_pos -= 15
                else:
                    c.setFont("Helvetica", 11)
                    c.drawString(margin + 10, y_pos, "Aucune donnée de visiteurs enregistrée")
                    y_pos -= 15
                    
            except Exception as visitor_error:
                logger.error(f"Erreur lors de la récupération des visiteurs: {visitor_error}")
                c.setFont("Helvetica", 11)
                c.drawString(margin + 10, y_pos, "Erreur dans les données de visiteurs")
                y_pos -= 15
            
            y_pos -= 10
            
            # Section présence individuelle (si espace disponible)
            if y_pos > margin + 80:  # Vérifier s'il y a assez d'espace
                try:
                    attendances = Attendance.objects.filter(event=event, attended=True).select_related('member__user')[:15]
                    
                    if attendances:
                        c.setFont("Helvetica-Bold", 14)
                        c.drawString(margin, y_pos, "PRÉSENCE INDIVIDUELLE (15 premiers)")
                        y_pos -= 20
                        
                        c.setFont("Helvetica", 10)
                        for attendance in attendances:
                            if y_pos < margin + 30:  # Nouvelle page si nécessaire
                                c.showPage()
                                y_pos = height - margin - 30
                                c.setFont("Helvetica", 10)
                            
                            member = attendance.member
                            if member and member.user:
                                name_parts = [
                                    member.user.first_name or '',
                                    member.post_name or '',
                                    member.user.last_name or ''
                                ]
                                name = ' '.join(filter(None, name_parts))
                                # Limiter la longueur pour éviter les dépassements
                                if len(name) > 45:
                                    name = name[:42] + "..."
                                c.drawString(margin + 10, y_pos, name)
                                y_pos -= 12
                
                except Exception as attendance_error:
                    logger.error(f"Erreur lors de la récupération des présences: {attendance_error}")
                    # Ne pas afficher d'erreur pour éviter de surcharger le PDF
            
            # Pied de page
            c.setFont("Helvetica-Oblique", 9)
            c.setFillColor(grey)
            footer_text = f"Généré le {timezone.now().strftime('%d/%m/%Y %H:%M')} - Consolation et Paix Divine"
            c.drawString(margin, margin - 10, footer_text)
            
            # Finaliser le PDF
            c.save()
            
            # Récupérer les données du PDF
            pdf_data = buf.getvalue()
            buf.close()
            
            # Vérifier que le PDF n'est pas vide
            if len(pdf_data) < 1000:  # Un PDF valide doit avoir plus de 1KB
                raise ValueError("PDF généré trop petit, probablement corrompu")
            
            # Créer la réponse HTTP
            filename = f"rapport_pointage_{event.id}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            response['Content-Length'] = len(pdf_data)

            logger.info(f"PDF généré avec succès: {filename} ({len(pdf_data)} bytes)")
            return response

        except Exception as e:
            logger.error(f"Erreur lors de la génération PDF: {str(e)}", exc_info=True)
            return Response(
                {'error': f"Erreur lors de la génération du rapport: {str(e)}"},
                status=500,
            )

    def _attendance_report_xlsx(self, request, event):
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.worksheet.table import Table, TableStyleInfo
        except ModuleNotFoundError:
            return Response(
                {
                    'detail': "Le module 'openpyxl' n'est pas installé sur le serveur. Installez-le puis relancez le backend.",
                },
                status=500,
            )

        wb = Workbook()
        ws = wb.active
        ws.title = 'Résumé'

        header_fill = PatternFill('solid', fgColor='1F4E79')
        header_font = Font(bold=True, color='FFFFFF')
        section_fill = PatternFill('solid', fgColor='0F766E')
        section_font = Font(bold=True, color='FFFFFF')
        thin = Side(style='thin', color='D9D9D9')
        border = Border(left=thin, right=thin, top=thin, bottom=thin)
        wrap_left = Alignment(horizontal='left', vertical='top', wrap_text=True)
        center = Alignment(horizontal='center', vertical='center', wrap_text=True)

        def add_title(row, text):
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
            c = ws.cell(row=row, column=1, value=text)
            c.fill = header_fill
            c.font = header_font
            c.alignment = center
            for col in range(1, 5):
                ws.cell(row=row, column=col).border = border
            ws.row_dimensions[row].height = 26

        def add_section(row, text):
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
            c = ws.cell(row=row, column=1, value=text)
            c.fill = section_fill
            c.font = section_font
            c.alignment = center
            for col in range(1, 5):
                ws.cell(row=row, column=col).border = border
            ws.row_dimensions[row].height = 20

        def add_kv(row, key, value):
            ws.cell(row=row, column=1, value=str(key or '')).font = Font(bold=True)
            ws.cell(row=row, column=2, value=value)
            ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=4)
            for col in range(1, 5):
                c = ws.cell(row=row, column=col)
                c.border = border
                c.alignment = wrap_left
            return row + 1

        r = 1
        add_title(r, 'RAPPORT DE POINTAGE')
        r += 2

        r = add_kv(r, 'Événement', getattr(event, 'title', None) or '—')
        r = add_kv(r, 'Date', getattr(event, 'date', None))
        r = add_kv(r, 'Heure', getattr(event, 'time', None))
        r = add_kv(r, 'Lieu', getattr(event, 'location', None) or '—')

        # Formats date/heure (col B)
        ws.cell(row=r - 3, column=2).number_format = 'dd/mm/yyyy'
        ws.cell(row=r - 2, column=2).number_format = 'hh:mm'

        r += 1
        add_section(r, 'PARTICIPANTS')
        r += 1

        ws.append(['Catégorie', 'Nombre', '', ''])
        for cidx in range(1, 5):
            c = ws.cell(row=r, column=cidx)
            c.font = Font(bold=True)
            c.alignment = center
            c.border = border
        r += 1

        attendance_agg = EventAttendanceAggregate.objects.filter(event=event).first()
        male_adults = (attendance_agg.male_adults or 0) if attendance_agg else 0
        female_adults = (attendance_agg.female_adults or 0) if attendance_agg else 0
        male_children = (attendance_agg.male_children or 0) if attendance_agg else 0
        female_children = (attendance_agg.female_children or 0) if attendance_agg else 0
        total_participants = male_adults + female_adults + male_children + female_children

        participant_rows = [
            ('Total participants', total_participants),
            ('Hommes adultes', male_adults),
            ('Femmes adultes', female_adults),
            ('Garçons', male_children),
            ('Filles', female_children),
        ]
        start_participants = r
        for label, val in participant_rows:
            ws.append([label, int(val or 0), '', ''])
            for cidx in range(1, 5):
                cell = ws.cell(row=r, column=cidx)
                cell.border = border
                cell.alignment = wrap_left if cidx == 1 else center
            r += 1
        end_participants = r - 1

        # Make it a table for better Excel formatting
        part_ref = f"A{start_participants - 1}:B{end_participants}"
        table = Table(displayName='ParticipantsTable', ref=part_ref)
        table.tableStyleInfo = TableStyleInfo(name='TableStyleMedium9', showRowStripes=True)
        ws.add_table(table)

        r += 1
        add_section(r, 'VISITEURS')
        r += 1

        ws.append(['Catégorie', 'Nombre', '', ''])
        for cidx in range(1, 5):
            c = ws.cell(row=r, column=cidx)
            c.font = Font(bold=True)
            c.alignment = center
            c.border = border
        r += 1

        visitor_agg = EventVisitorAggregate.objects.filter(event=event).first()
        male_visitors = (visitor_agg.male_visitors or 0) if visitor_agg else 0
        female_visitors = (visitor_agg.female_visitors or 0) if visitor_agg else 0
        total_visitors = male_visitors + female_visitors
        visitor_rows = [
            ('Total visiteurs', total_visitors),
            ('Hommes visiteurs', male_visitors),
            ('Femmes visiteuses', female_visitors),
        ]
        start_visitors = r
        for label, val in visitor_rows:
            ws.append([label, int(val or 0), '', ''])
            for cidx in range(1, 5):
                cell = ws.cell(row=r, column=cidx)
                cell.border = border
                cell.alignment = wrap_left if cidx == 1 else center
            r += 1
        end_visitors = r - 1

        vis_ref = f"A{start_visitors - 1}:B{end_visitors}"
        vtable = Table(displayName='VisitorsTable', ref=vis_ref)
        vtable.tableStyleInfo = TableStyleInfo(name='TableStyleMedium9', showRowStripes=True)
        ws.add_table(vtable)

        # Column widths
        ws.column_dimensions['A'].width = 28
        ws.column_dimensions['B'].width = 14
        ws.column_dimensions['C'].width = 2
        ws.column_dimensions['D'].width = 38

        ws.freeze_panes = 'A3'
        ws.sheet_view.showGridLines = False

        # Second sheet: presence list
        ws2 = wb.create_sheet('Présences')
        headers2 = ['Prénom', 'Postnom', 'Nom', 'Téléphone']
        ws2.append(headers2)
        ws2.row_dimensions[1].height = 24
        for col_idx, _ in enumerate(headers2, start=1):
            c = ws2.cell(row=1, column=col_idx)
            c.fill = header_fill
            c.font = header_font
            c.alignment = center
            c.border = border

        attendances = Attendance.objects.filter(event=event, attended=True).select_related('member__user')
        for attendance in attendances:
            member = attendance.member
            u = member.user if member else None
            ws2.append([
                getattr(u, 'first_name', '') if u else '',
                getattr(member, 'post_name', '') if member else '',
                getattr(u, 'last_name', '') if u else '',
                getattr(u, 'phone', '') if u else '',
            ])

        last_row2 = ws2.max_row
        last_col2 = ws2.max_column
        for row in ws2.iter_rows(min_row=2, max_row=last_row2, min_col=1, max_col=last_col2):
            for cell in row:
                cell.alignment = wrap_left
                cell.border = border

        ws2.freeze_panes = 'A2'
        ws2.sheet_view.showGridLines = False

        if last_row2 >= 2:
            table_ref2 = f"A1:{get_column_letter(last_col2)}{last_row2}"
            t2 = Table(displayName='AttendanceList', ref=table_ref2)
            t2.tableStyleInfo = TableStyleInfo(name='TableStyleMedium9', showRowStripes=True)
            ws2.add_table(t2)

        for col_idx in range(1, last_col2 + 1):
            max_len = 0
            for cell in ws2[get_column_letter(col_idx)]:
                v = cell.value
                if v is None:
                    continue
                s = str(v)
                if len(s) > max_len:
                    max_len = len(s)
            ws2.column_dimensions[get_column_letter(col_idx)].width = min(max(12, max_len + 2), 40)

        out = io.BytesIO()
        wb.save(out)
        out.seek(0)

        filename = f"rapport_pointage_{event.id}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        response = HttpResponse(
            out.getvalue(),
            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        )
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

    def _attendance_report_csv(self, request, event):
        import csv
        from django.utils.encoding import smart_str
        
        response = HttpResponse(content_type='text/csv; charset=utf-8')
        filename = f"rapport_pointage_{event.id}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.csv"
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        
        writer = csv.writer(response)
        
        try:
            # En-tête CSV
            writer.writerow(['RAPPORT DE POINTAGE'])
            writer.writerow(['Événement', event.title or 'Non spécifié'])
            
            date_str = event.date.strftime('%d/%m/%Y') if event.date else 'Non spécifiée'
            writer.writerow(['Date', date_str])
            
            if event.time:
                time_str = event.time.strftime('%H:%M') if event.time else ''
                writer.writerow(['Heure', time_str])
            
            if event.location:
                writer.writerow(['Lieu', event.location])
            
            writer.writerow([])
            
            # Participants
            writer.writerow(['PARTICIPANTS'])
            writer.writerow(['Catégorie', 'Nombre'])
            
            attendance_agg = EventAttendanceAggregate.objects.filter(event=event).first()
            
            if attendance_agg:
                male_adults = attendance_agg.male_adults or 0
                female_adults = attendance_agg.female_adults or 0
                male_children = attendance_agg.male_children or 0
                female_children = attendance_agg.female_children or 0
                
                total_participants = male_adults + female_adults + male_children + female_children
                writer.writerow(['Total participants', total_participants])
                writer.writerow(['Hommes adultes', male_adults])
                writer.writerow(['Femmes adultes', female_adults])
                writer.writerow(['Garçons', male_children])
                writer.writerow(['Filles', female_children])
            else:
                writer.writerow(['Total participants', 0])
                writer.writerow(['Hommes adultes', 0])
                writer.writerow(['Femmes adultes', 0])
                writer.writerow(['Garçons', 0])
                writer.writerow(['Filles', 0])
            
            writer.writerow([])
            
            # Visiteurs
            writer.writerow(['VISITEURS'])
            writer.writerow(['Catégorie', 'Nombre'])
            
            visitor_agg = EventVisitorAggregate.objects.filter(event=event).first()
            
            if visitor_agg:
                male_visitors = visitor_agg.male_visitors or 0
                female_visitors = visitor_agg.female_visitors or 0
                
                total_visitors = male_visitors + female_visitors
                writer.writerow(['Total visiteurs', total_visitors])
                writer.writerow(['Hommes visiteurs', male_visitors])
                writer.writerow(['Femmes visiteuses', female_visitors])
            else:
                writer.writerow(['Total visiteurs', 0])
                writer.writerow(['Hommes visiteurs', 0])
                writer.writerow(['Femmes visiteuses', 0])
            
            writer.writerow([])
            
            # Présence individuelle
            attendances = Attendance.objects.filter(event=event, attended=True).select_related('member__user')
            if attendances.exists():
                writer.writerow(['PRÉSENCE INDIVIDUELLE'])
                writer.writerow(['Nom', 'Postnom', 'Nom de famille', 'Téléphone'])
                
                for attendance in attendances:
                    member = attendance.member
                    if member and member.user:
                        first_name = member.user.first_name or ''
                        post_name = member.post_name or ''
                        last_name = member.user.last_name or ''
                        phone = member.user.phone or ''
                        writer.writerow([first_name, post_name, last_name, phone])
            else:
                writer.writerow(['PRÉSENCE INDIVIDUELLE'])
                writer.writerow(['Aucune présence enregistrée'])
            
        except Exception as e:
            writer.writerow(['ERREUR'])
            writer.writerow(['Message', f'Erreur lors de la génération: {str(e)}'])
        
        return response

    def _event_program_pdf(self, request, event):
        width, height = A4
        margin = 18 * mm

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)

        def safe(s):
            return (s or '').strip()

        title = safe(event.title) or 'Programme'
        desc = safe(getattr(event, 'description', None))
        location = safe(getattr(event, 'location', None))

        event_type = safe(getattr(event, 'event_type', None))
        duration_type = safe(getattr(event, 'duration_type', None))
        date_str = event.date.strftime('%Y-%m-%d') if getattr(event, 'date', None) else ''
        time_str = event.time.strftime('%H:%M') if getattr(event, 'time', None) else ''

        def label_for_choice(choices, value):
            if not value:
                return ''
            for v, lbl in choices:
                if v == value:
                    return lbl
            return value

        event_type_lbl = label_for_choice(getattr(Event, 'EVENT_TYPE_CHOICES', []), event_type)
        duration_obj = ActivityDuration.objects.filter(code=duration_type).only('label').first()
        duration_lbl = getattr(duration_obj, 'label', None) or label_for_choice(getattr(Event, 'DURATION_CHOICES', []), duration_type)

        frontend_base = getattr(settings, 'FRONTEND_BASE_URL', 'http://localhost:3000')
        public_link = None
        if getattr(event, 'share_slug', None):
            public_link = f"{frontend_base.rstrip('/')}/p/{event.share_slug}"

        cert = None
        verify_url = None
        try:
            payload = {
                'event_id': event.id,
                'date': getattr(event, 'date', None).isoformat() if getattr(event, 'date', None) else None,
                'title': getattr(event, 'title', None),
                'share_slug': getattr(event, 'share_slug', None),
                'generated_at': timezone.now().isoformat(),
            }
            cert = _create_report_certificate('programme', payload, getattr(request, 'user', None))
            verify_url = request.build_absolute_uri(f"/api/reports/verify/?code={cert.code}")
        except Exception:
            cert = None
            verify_url = None

        # Background
        is_published = bool(getattr(event, 'is_published', False))
        meta = FinancialTransactionViewSet()._pdf_brand_layout(
            c,
            doc_title='Programme',
            doc_subtitle='Document officiel — diffusion interne / publique',
            badge_text='PUBLIÉ' if is_published else 'BROUILLON',
            badge_rgb=(0.05, 0.65, 0.45) if is_published else (0.45, 0.50, 0.60),
        )
        width = meta['width']
        height = meta['height']
        margin = meta['margin']
        card_top = meta['card_top']

        # Header
        header_h = 32 * mm

        # Poster (optional) as banner
        poster_drawn = False
        poster_h = 46 * mm
        poster_w = width - 2 * margin
        poster_top = card_top - 8 * mm
        poster_y = poster_top - poster_h
        if poster_y < 22 * mm:
            poster_y = 22 * mm
            poster_h = max(0, float(poster_top) - float(poster_y))
        try:
            if getattr(event, 'poster_image', None) and getattr(event.poster_image, 'path', None):
                img = ImageReader(event.poster_image.path)
                c.drawImage(img, margin, poster_y, width=poster_w, height=poster_h, preserveAspectRatio=True, anchor='c', mask='auto')
                poster_drawn = True
        except Exception:
            poster_drawn = False

        if not poster_drawn:
            c.setFillColorRGB(0.93, 0.95, 0.99)
            c.roundRect(margin, poster_y, poster_w, poster_h, 10, fill=1, stroke=0)
            c.setFillColorRGB(0.20, 0.35, 0.75)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(margin + 10 * mm, poster_y + poster_h - 14 * mm, 'Affiche non disponible')

        # Main card
        x = margin + 10 * mm
        y = poster_y - 12 * mm

        # Title
        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 18)
        c.drawString(x, y, title)

        # Status badge
        y -= 2 * mm

        y -= 10 * mm

        # Key facts row
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica', 10)
        facts_left = f"{event_type_lbl or event_type}  •  {duration_lbl or duration_type}"
        facts_mid = f"{date_str}  •  {time_str}".strip(' •')
        c.drawString(x, y, facts_left.strip())
        c.drawRightString(width - margin - 10 * mm, y, facts_mid)

        y -= 8 * mm
        if location:
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 10)
            c.drawString(x, y, 'Lieu')
            c.setFont('Helvetica', 10)
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.drawString(x + 16 * mm, y, location)
            y -= 8 * mm

        # Intervenants
        people = [
            ('Modérateur', safe(getattr(event, 'moderator', None))),
            ('Prédicateur', safe(getattr(event, 'preacher', None))),
            ('Chorale', safe(getattr(event, 'choir', None))),
            ('Protocole & Diaconat', safe(getattr(event, 'protocol_team', None))),
            ('Sonorisation/Technique/Media', safe(getattr(event, 'tech_team', None))),
            ('Communicateur', safe(getattr(event, 'communicator', None))),
        ]
        people = [(k, v) for (k, v) in people if v]

        if people:
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Intervenants')
            y -= 6 * mm

            c.setFont('Helvetica', 10)
            c.setFillColorRGB(0.20, 0.24, 0.30)
            col_w = (width - 2 * margin - 20 * mm) / 2
            left_x = x
            right_x = x + col_w + 10 * mm
            row_y = y
            for i, (k, v) in enumerate(people):
                col_x = left_x if i % 2 == 0 else right_x
                if i % 2 == 0 and i > 0:
                    row_y -= 6 * mm
                c.setFillColorRGB(0.40, 0.45, 0.55)
                c.setFont('Helvetica-Bold', 9)
                c.drawString(col_x, row_y, k)
                c.setFillColorRGB(0.20, 0.24, 0.30)
                c.setFont('Helvetica', 10)
                c.drawString(col_x + 34 * mm, row_y, v)

            y = row_y - 10 * mm
        else:
            y -= 2 * mm

        # Description
        if desc:
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 11)
            c.drawString(x, y, 'Description')
            y -= 6 * mm

            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            max_w = width - margin - x - 10 * mm
            words = desc.replace('\r', '').split('\n')
            lines = []
            for paragraph in words:
                paragraph = paragraph.strip()
                if not paragraph:
                    lines.append('')
                    continue
                cur = ''
                for w in paragraph.split(' '):
                    candidate = (cur + ' ' + w).strip()
                    if c.stringWidth(candidate, 'Helvetica', 10) <= max_w:
                        cur = candidate
                    else:
                        if cur:
                            lines.append(cur)
                        cur = w
                if cur:
                    lines.append(cur)
            for ln in lines[:18]:
                c.drawString(x, y, ln)
                y -= 5.2 * mm
            y -= 2 * mm

        if cert and verify_url:
            try:
                _draw_authenticity_qr(c, verify_url, cert.code, x, width, margin, qr_size_mm=26, qr_y_mm=34, text_top_mm=46)
            except Exception:
                pass

        if public_link:
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 10)
            c.drawString(x, 28 * mm, 'Lien public')
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 9)
            c.drawString(x, 23 * mm, public_link[:110])

        # Footer
        c.setFillColorRGB(0.45, 0.50, 0.60)
        c.setFont('Helvetica', 8)
        gen_at = timezone.now().strftime('%Y-%m-%d %H:%M')
        c.drawRightString(width - margin, 10 * mm, f"Généré le {gen_at}")

        c.showPage()
        c.save()
        pdf = buf.getvalue()
        buf.close()

        if cert:
            try:
                cert.pdf_sha256 = hashlib.sha256(pdf).hexdigest()
                cert.save(update_fields=['pdf_sha256', 'updated_at'])
            except Exception:
                pass

        return pdf

    @action(detail=True, methods=['get'], url_path='pdf')
    def pdf(self, request, pk=None):
        event = self.get_object()
        self._ensure_share_slug(event)
        if not getattr(event, 'is_published', False) and not (request.user and request.user.is_authenticated):
            return Response({'detail': 'non autorisé'}, status=403)

        pdf_bytes = self._event_program_pdf(request, event)
        title = (event.title or 'programme').strip()
        title = re.sub(r'[^a-zA-Z0-9\-_ ]+', '', title).strip().replace(' ', '_')
        date_str = event.date.strftime('%Y%m%d') if getattr(event, 'date', None) else timezone.now().strftime('%Y%m%d')
        filename = f"programme_{date_str}_{title or 'event'}.pdf"
        resp = HttpResponse(pdf_bytes, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp

    @action(detail=True, methods=['get', 'post'], url_path='comments')
    def comments(self, request, pk=None):
        event = self.get_object()
        if request.method == 'GET':
            qs = EventComment.objects.filter(event=event).order_by('-created_at')
            return Response(EventCommentSerializer(qs, many=True).data)

        serializer = EventCommentSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        author = (f"{request.user.first_name} {request.user.last_name}").strip() if request.user else ''
        author = author or getattr(request.user, 'username', None) or serializer.validated_data.get('author_name') or '—'
        serializer.save(event=event, author_name=author)
        return Response(serializer.data, status=201)

    @action(detail=False, methods=['get'], url_path=r'public/(?P<slug>[^/.]+)')
    def public(self, request, slug=None):
        event = Event.objects.filter(share_slug=slug, is_published=True).order_by('-date', '-time').first()
        if not event:
            return Response({'detail': 'événement introuvable'}, status=404)

        data = EventSerializer(event).data
        comments = EventComment.objects.filter(event=event).order_by('-created_at')[:20]
        data['comments'] = EventCommentSerializer(comments, many=True).data
        return Response(data)

    @action(detail=False, methods=['post'], url_path=r'public/(?P<slug>[^/.]+)/comment')
    def public_comment(self, request, slug=None):
        event = Event.objects.filter(share_slug=slug, is_published=True).order_by('-date', '-time').first()
        if not event:
            return Response({'detail': 'événement introuvable'}, status=404)

        serializer = EventCommentSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        author = serializer.validated_data.get('author_name') or 'Anonyme'
        serializer.save(event=event, author_name=author)
        return Response(serializer.data, status=201)


class AttendanceViewSet(viewsets.ModelViewSet):
    queryset = Attendance.objects.all().order_by('-id')
    serializer_class = AttendanceSerializer
    permission_classes = [IsSecretaryOrAdmin]

    def get_queryset(self):
        qs = super().get_queryset()
        event_id = (self.request.query_params.get('event') or '').strip()
        if event_id:
            try:
                qs = qs.filter(event_id=int(event_id))
            except (TypeError, ValueError):
                return qs.none()

        attended = (self.request.query_params.get('attended') or '').strip().lower()
        if attended in {'1', 'true', 'yes'}:
            qs = qs.filter(attended=True)
        elif attended in {'0', 'false', 'no'}:
            qs = qs.filter(attended=False)

        return qs


class FinancialCategoryViewSet(viewsets.ModelViewSet):
    queryset = FinancialCategory.objects.all().order_by('-id')
    serializer_class = FinancialCategorySerializer
    permission_classes = [IsTreasurerOrAdmin]


class FinancialTransactionViewSet(viewsets.ModelViewSet):
    queryset = FinancialTransaction.objects.all().order_by('-date', '-id')
    serializer_class = FinancialTransactionSerializer
    permission_classes = [IsTreasurerOrAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def create(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().create(request, *args, **kwargs)
        ar = _create_approval_request(request, model='FinancialTransaction', action='create', payload=_safe_payload(request.data))
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def update(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().update(request, *args, **kwargs)
        tx = self.get_object()
        ar = _create_approval_request(
            request,
            model='FinancialTransaction',
            action='update',
            payload=_safe_payload(request.data),
            target_object_id=getattr(tx, 'id', None),
            object_repr=getattr(tx, 'document_number', None) or getattr(tx, 'receipt_code', None) or str(getattr(tx, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def partial_update(self, request, *args, **kwargs):
        return self.update(request, *args, **kwargs)

    def destroy(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().destroy(request, *args, **kwargs)
        tx = self.get_object()
        ar = _create_approval_request(
            request,
            model='FinancialTransaction',
            action='delete',
            payload=None,
            target_object_id=getattr(tx, 'id', None),
            object_repr=getattr(tx, 'document_number', None) or getattr(tx, 'receipt_code', None) or str(getattr(tx, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def get_queryset(self):
        qs = super().get_queryset()

        event_id = (self.request.query_params.get('event') or '').strip()
        if event_id:
            try:
                qs = qs.filter(event_id=int(event_id))
            except (TypeError, ValueError):
                qs = qs.none()

        direction = (self.request.query_params.get('direction') or '').strip().lower()
        if direction in {'in', 'out'}:
            qs = qs.filter(direction=direction)

        tx_type = (self.request.query_params.get('transaction_type') or '').strip().lower()
        if tx_type:
            qs = qs.filter(transaction_type=tx_type)

        currency = (self.request.query_params.get('currency') or '').strip().upper()
        if currency:
            qs = qs.filter(currency=currency)

        start = (self.request.query_params.get('start') or '').strip()
        end = (self.request.query_params.get('end') or '').strip()
        if start and end:
            qs = qs.filter(date__range=[start, end])
        elif start:
            qs = qs.filter(date__gte=start)
        elif end:
            qs = qs.filter(date__lte=end)

        q = (self.request.query_params.get('q') or '').strip()
        if q:
            qs = qs.filter(
                Q(document_number__icontains=q)
                | Q(receipt_code__icontains=q)
                | Q(reference_number__icontains=q)
                | Q(donor_name__icontains=q)
                | Q(donor_email__icontains=q)
                | Q(recipient_name__icontains=q)
                | Q(recipient_email__icontains=q)
                | Q(recipient_phone__icontains=q)
            )

        return qs

    @action(detail=False, methods=['get'], url_path='export')
    def export(self, request):
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.worksheet.table import Table, TableStyleInfo
        except ModuleNotFoundError:
            return Response(
                {
                    'detail': "Le module 'openpyxl' n'est pas installé sur le serveur. Installez-le puis relancez le backend.",
                },
                status=500,
            )

        wb = Workbook()
        ws = wb.active
        ws.title = 'Transactions'

        headers = [
            'Document',
            'Sens',
            'Type',
            'Montant',
            'Devise',
            'Date',
            'Donateur',
            'Email donateur',
            'Bénéficiaire',
            'Email bénéficiaire',
            'Téléphone bénéficiaire',
            'Mode paiement',
            'Référence',
            'Description',
            'Caissier',
            'Créé par',
            'Créé le',
        ]
        ws.append(headers)

        header_fill = PatternFill('solid', fgColor='7C2D12')
        header_font = Font(bold=True, color='FFFFFF')
        header_alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        thin = Side(style='thin', color='D9D9D9')
        border = Border(left=thin, right=thin, top=thin, bottom=thin)

        ws.row_dimensions[1].height = 28
        for col_idx, _ in enumerate(headers, start=1):
            c = ws.cell(row=1, column=col_idx)
            c.fill = header_fill
            c.font = header_font
            c.alignment = header_alignment
            c.border = border

        qs = self.get_queryset().select_related('cashier', 'created_by')
        for tx in qs:
            cashier_name = None
            if tx.cashier:
                cashier_name = (f"{tx.cashier.first_name} {tx.cashier.last_name}").strip() or tx.cashier.username
            created_by_name = None
            if tx.created_by:
                created_by_name = (f"{tx.created_by.first_name} {tx.created_by.last_name}").strip() or tx.created_by.username

            created_at = getattr(tx, 'created_at', None)
            if created_at:
                try:
                    created_at = timezone.localtime(created_at).replace(tzinfo=None)
                except Exception:
                    created_at = None

            ws.append([
                tx.document_number or tx.receipt_code or '',
                tx.direction or '',
                tx.transaction_type or '',
                float(tx.amount) if tx.amount is not None else None,
                (tx.currency or '').upper(),
                getattr(tx, 'date', None),
                tx.donor_name or '',
                tx.donor_email or '',
                tx.recipient_name or '',
                tx.recipient_email or '',
                tx.recipient_phone or '',
                tx.payment_method or '',
                tx.reference_number or '',
                tx.description or '',
                cashier_name or '',
                created_by_name or '',
                created_at,
            ])

        last_row = ws.max_row
        last_col = ws.max_column

        body_alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
        money_alignment = Alignment(horizontal='right', vertical='top')

        for row in ws.iter_rows(min_row=2, max_row=last_row, min_col=1, max_col=last_col):
            for cell in row:
                cell.alignment = body_alignment
                cell.border = border

        for r in range(2, last_row + 1):
            ws.cell(row=r, column=4).number_format = '#,##0.00'
            ws.cell(row=r, column=4).alignment = money_alignment
            ws.cell(row=r, column=6).number_format = 'dd/mm/yyyy'
            ws.cell(row=r, column=17).number_format = 'dd/mm/yyyy hh:mm'

        ws.freeze_panes = 'A2'

        if last_row >= 2:
            table_ref = f"A1:{get_column_letter(last_col)}{last_row}"
            table = Table(displayName='FinancialTransactionsTable', ref=table_ref)
            style = TableStyleInfo(
                name='TableStyleMedium9',
                showFirstColumn=False,
                showLastColumn=False,
                showRowStripes=True,
                showColumnStripes=False,
            )
            table.tableStyleInfo = style
            ws.add_table(table)

        for col_idx in range(1, last_col + 1):
            max_len = 0
            for cell in ws[get_column_letter(col_idx)]:
                v = cell.value
                if v is None:
                    continue
                if hasattr(v, 'strftime'):
                    s = v.strftime('%d/%m/%Y')
                else:
                    s = str(v)
                if len(s) > max_len:
                    max_len = len(s)
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max(10, max_len + 2), 48)

        ws.sheet_view.showGridLines = False
        ws.print_title_rows = '1:1'
        ws.page_setup.orientation = 'landscape'
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0

        out = io.BytesIO()
        wb.save(out)
        out.seek(0)

        filename = f"financial_transactions_{timezone.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        resp = HttpResponse(
            out.getvalue(),
            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        )
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp

    def get_permissions(self):
        if getattr(self, 'action', None) in {'verify_receipt', 'verify_document'}:
            return [AllowAny()]
        return super().get_permissions()

    def _ensure_document_number(self, tx):
        if tx.document_number:
            return tx.document_number

        year = tx.date.year if tx.date else timezone.now().year
        prefix = 'RC' if tx.direction == 'in' else 'BS'

        seq, _ = FinancialDocumentSequence.objects.select_for_update().get_or_create(
            prefix=prefix,
            year=year,
            defaults={'last_number': 0},
        )
        seq.last_number = (seq.last_number or 0) + 1
        seq.save(update_fields=['last_number'])

        tx.document_number = f"{prefix}-{year}-{seq.last_number:06d}"
        tx.save(update_fields=['document_number'])
        return tx.document_number

    def _ensure_receipt_code(self, tx):
        if tx.receipt_code:
            return tx.receipt_code

        doc = self._ensure_document_number(tx)
        tx.receipt_code = doc
        tx.save(update_fields=['receipt_code'])
        return tx.receipt_code

    def _pdf_get_brand_logo_path(self, kind='header'):
        base_dir = getattr(settings, 'BASE_DIR', None) or ''
        base_dir = str(base_dir)
        brand_dirs = []
        if base_dir:
            brand_dirs.append(os.path.join(base_dir, 'brand'))
            try:
                brand_dirs.append(os.path.join(os.path.dirname(base_dir), 'brand'))
            except Exception:
                pass

        candidates = []
        if kind == 'watermark':
            for bd in brand_dirs:
                candidates.extend([
                    os.path.join(bd, 'logo-pdf-watermark.png'),
                    os.path.join(bd, 'logo-pdf-watermark.jpg'),
                    os.path.join(bd, 'logo-pdf-watermark.jpeg'),
                ])
        else:
            for bd in brand_dirs:
                candidates.extend([
                    os.path.join(bd, 'logo-pdf-header.png'),
                    os.path.join(bd, 'logo-pdf-header.jpg'),
                    os.path.join(bd, 'logo-pdf-header.jpeg'),
                ])

        for bd in brand_dirs:
            candidates.extend([
                os.path.join(bd, 'logo-pdf.png'),
                os.path.join(bd, 'logo-pdf.jpg'),
                os.path.join(bd, 'logo-pdf.jpeg'),
                os.path.join(bd, 'logo.png'),
                os.path.join(bd, 'logo.jpg'),
                os.path.join(bd, 'logo.jpeg'),
                os.path.join(bd, 'logo-dark.png'),
                os.path.join(bd, 'logo-light.png'),
            ])
        for p in candidates:
            try:
                if p and os.path.exists(p):
                    return p
            except Exception:
                continue
        return None

    def _pdf_draw_brand_logo(self, c, x, y, box_w, box_h, kind='header'):
        path = self._pdf_get_brand_logo_path(kind=kind)
        if not path:
            return False
        try:
            img = ImageReader(path)
            iw, ih = img.getSize()
            mask = 'auto'
            try:
                if os.path.basename(str(path)).lower() == 'logo.png':
                    mask = [0, 18, 0, 18, 0, 18]
            except Exception:
                mask = 'auto'
            if iw and ih:
                scale = min(float(box_w) / float(iw), float(box_h) / float(ih))
                w = float(iw) * scale
                h = float(ih) * scale
                c.drawImage(img, x, y + (float(box_h) - h) / 2.0, w, h, preserveAspectRatio=True, mask=mask)
            else:
                c.drawImage(img, x, y, box_w, box_h, preserveAspectRatio=True, mask=mask)
            return True
        except Exception:
            return False

    def _pdf_draw_brand_watermark(self, c, x0, y0, x1, y1, opacity=0.06):
        path = self._pdf_get_brand_logo_path(kind='watermark')
        if not path:
            return

        try:
            img = ImageReader(path)
            iw, ih = img.getSize()
            if not iw or not ih:
                return

            mask = 'auto'
            try:
                if os.path.basename(str(path)).lower() == 'logo.png':
                    mask = [0, 18, 0, 18, 0, 18]
            except Exception:
                mask = 'auto'

            bw = float(x1) - float(x0)
            bh = float(y1) - float(y0)
            max_w = bw * 0.58
            max_h = bh * 0.58
            scale = min(max_w / float(iw), max_h / float(ih))
            w = float(iw) * scale
            h = float(ih) * scale
            x = float(x0) + (bw - w) / 2.0
            y = float(y0) + (bh - h) / 2.0

            c.saveState()
            try:
                c.setFillAlpha(opacity)
                c.setStrokeAlpha(opacity)
            except Exception:
                pass
            c.drawImage(img, x, y, w, h, preserveAspectRatio=True, mask=mask)
        except Exception:
            return
        finally:
            try:
                c.restoreState()
            except Exception:
                pass

    def _pdf_brand_layout(self, c, doc_title, doc_subtitle, badge_text=None, badge_rgb=None):
        width, height = A4
        margin = 18 * mm

        c.setFillColorRGB(0.98, 0.985, 0.995)
        c.rect(0, 0, width, height, fill=1, stroke=0)

        header_h = 32 * mm
        c.setFillColorRGB(0.06, 0.18, 0.45)
        c.rect(0, height - header_h, width, header_h, fill=1, stroke=0)

        logo_box = 16 * mm
        logo_y = height - header_h + (header_h - logo_box) / 2.0
        has_logo = self._pdf_draw_brand_logo(c, margin, logo_y, logo_box, logo_box, kind='header')
        text_left = margin + (logo_box + 4 * mm if has_logo else 0)

        c.setFillColorRGB(1, 1, 1)
        c.setFont('Helvetica-Bold', 15)
        c.drawString(text_left, height - 13 * mm, 'Consolation et Paix Divine')

        c.setFillColorRGB(0.92, 0.95, 1)
        c.setFont('Helvetica', 9)
        c.drawString(text_left, height - 20 * mm, (doc_subtitle or '').strip())

        c.setFillColorRGB(1, 1, 1)
        c.setFont('Helvetica-Bold', 14)
        c.drawRightString(width - margin, height - 13 * mm, (doc_title or '').strip())

        if badge_text:
            bw = 26 * mm
            bh = 7 * mm
            bx = width - margin - bw
            by = height - header_h + 6.5 * mm
            if badge_rgb:
                c.setFillColorRGB(*badge_rgb)
            else:
                c.setFillColorRGB(0.45, 0.50, 0.60)
            c.roundRect(bx, by, bw, bh, 7, fill=1, stroke=0)
            c.setFillColorRGB(1, 1, 1)
            c.setFont('Helvetica-Bold', 9)
            c.drawCentredString(bx + bw / 2, by + 2.2 * mm, str(badge_text))

        card_top = height - header_h - 12 * mm
        card_h = card_top - 18 * mm
        c.setFillColorRGB(1, 1, 1)
        c.roundRect(margin, 18 * mm, width - 2 * margin, card_h, 12, fill=1, stroke=0)

        self._pdf_draw_brand_watermark(c, margin, 18 * mm, width - margin, card_top)

        return {
            'width': width,
            'height': height,
            'margin': margin,
            'card_top': card_top,
            'card_h': card_h,
            'x': margin + 10 * mm,
            'y_start': card_top - 12 * mm,
        }

    def _pdf_half_brand_layout(self, c, y0, y1, doc_title, doc_subtitle, badge_text=None, badge_rgb=None):
        width, _ = A4
        margin = 14 * mm

        header_h = 18 * mm
        c.setFillColorRGB(0.06, 0.18, 0.45)
        c.rect(0, y1 - header_h, width, header_h, fill=1, stroke=0)

        logo_box = 12 * mm
        logo_y = y1 - header_h + (header_h - logo_box) / 2.0
        has_logo = self._pdf_draw_brand_logo(c, margin, logo_y, logo_box, logo_box, kind='header')
        text_left = margin + (logo_box + 3 * mm if has_logo else 0)

        c.setFillColorRGB(1, 1, 1)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(text_left, y1 - 8 * mm, 'Consolation et Paix Divine')

        c.setFillColorRGB(0.92, 0.95, 1)
        c.setFont('Helvetica', 8)
        c.drawString(text_left, y1 - 13.2 * mm, (doc_subtitle or '').strip())

        c.setFillColorRGB(1, 1, 1)
        c.setFont('Helvetica-Bold', 12)
        c.drawRightString(width - margin, y1 - 8 * mm, (doc_title or '').strip())

        if badge_text:
            bw = 24 * mm
            bh = 6.5 * mm
            bx = width - margin - bw
            by = y1 - header_h + 5 * mm
            if badge_rgb:
                c.setFillColorRGB(*badge_rgb)
            else:
                c.setFillColorRGB(0.45, 0.50, 0.60)
            c.roundRect(bx, by, bw, bh, 7, fill=1, stroke=0)
            c.setFillColorRGB(1, 1, 1)
            c.setFont('Helvetica-Bold', 8)
            c.drawCentredString(bx + bw / 2, by + 2.0 * mm, str(badge_text))

        card_top = y1 - header_h - 6 * mm
        card_bottom = y0 + 10 * mm
        card_h = max(0, card_top - card_bottom)
        c.setFillColorRGB(1, 1, 1)
        c.roundRect(margin, card_bottom, width - 2 * margin, card_h, 10, fill=1, stroke=0)

        self._pdf_draw_brand_watermark(c, margin, card_bottom, width - margin, card_top)

        return {
            'width': width,
            'margin': margin,
            'x': margin + 8 * mm,
            'y_start': card_top - 8 * mm,
            'y0': y0,
            'y1': y1,
        }

    def _pdf_draw_signatures(self, c, x, y, block_w, left_label='Signature', right_label='Signature', left_name=None, right_name=None):
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 7.5)

        col_w = block_w / 2.0
        line_w = col_w - 6 * mm

        left_text = str(left_label or 'Signature')
        if left_name:
            left_text = f"{left_text}: {str(left_name)[:40]}"
        c.drawString(x, y + 2.5 * mm, left_text)
        c.line(x, y, x + line_w, y)

        cx = x + col_w
        right_text = str(right_label or 'Signature')
        if right_name:
            right_text = f"{right_text}: {str(right_name)[:40]}"
        c.drawString(cx, y + 2.5 * mm, right_text)
        c.line(cx, y, cx + line_w, y)

    def _build_receipt_pdf_bytes(self, tx, request):
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)

        meta = self._pdf_brand_layout(
            c,
            doc_title='REÇU',
            doc_subtitle='Finance — Document officiel',
            badge_text='ENTRÉE',
            badge_rgb=(0.10, 0.55, 0.35),
        )
        width = meta['width']
        margin = meta['margin']
        x = meta['x']
        y = meta['y_start']

        doc_no = tx.document_number or tx.receipt_code or ''

        donor = tx.donor_name or (tx.member.user.get_full_name() if tx.member and tx.member.user else None) or '—'
        cashier_name = None
        if tx.cashier:
            cashier_name = (f"{tx.cashier.first_name} {tx.cashier.last_name}").strip() or tx.cashier.username
        cashier_name = cashier_name or '—'

        created_by_name = None
        if tx.created_by:
            created_by_name = (f"{tx.created_by.first_name} {tx.created_by.last_name}").strip() or tx.created_by.username
        created_by_name = created_by_name or '—'

        verify_url = request.build_absolute_uri(f"/api/financial-transactions/verify-receipt/?code={tx.receipt_code}")

        def tx_type_label(value):
            try:
                for v, lbl in getattr(FinancialTransaction, 'TRANSACTION_TYPE_CHOICES', []):
                    if v == value:
                        return lbl
            except Exception:
                return value
            return value

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Informations')
        y -= 7 * mm

        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(x, y, 'Numéro')
        c.drawString(x + 55 * mm, y, 'Date')
        y -= 5.5 * mm
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 10)
        c.drawString(x, y, str(doc_no)[:40])
        c.drawString(x + 55 * mm, y, str(tx.date.isoformat())[:10])
        y -= 9 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Donneur')
        y -= 7 * mm
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 11)
        c.drawString(x, y, str(donor)[:80])
        y -= 7 * mm
        c.setFillColorRGB(0.45, 0.50, 0.60)
        c.setFont('Helvetica', 9)
        c.drawString(x, y, f"Email: {tx.donor_email or '—'}")
        y -= 10 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Montant')
        y -= 8 * mm
        c.setFillColorRGB(0.06, 0.18, 0.45)
        c.setFont('Helvetica-Bold', 22)
        c.drawString(x, y, f"{tx.amount} {tx.currency}")
        y -= 12 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Détails')
        y -= 7 * mm

        rows = [
            ('Type', tx_type_label(tx.transaction_type) or '—'),
            ('Mode de paiement', tx.payment_method or '—'),
            ('Référence', tx.reference_number or '—'),
            ('Motif/Note', tx.description or '—'),
            ('Fait par', created_by_name),
            ('Approuvé par', cashier_name),
        ]

        label_w = 38 * mm
        value_w = (width - margin - x - 10 * mm) - label_w
        for label, val in rows:
            if y < 32 * mm:
                c.showPage()
                meta = self._pdf_brand_layout(
                    c,
                    doc_title='REÇU',
                    doc_subtitle='Finance — Document officiel',
                    badge_text='ENTRÉE',
                    badge_rgb=(0.10, 0.55, 0.35),
                )
                width = meta['width']
                margin = meta['margin']
                x = meta['x']
                y = meta['y_start']
                c.setFillColorRGB(0.07, 0.10, 0.16)
                c.setFont('Helvetica-Bold', 12)
                c.drawString(x, y, 'Détails')
                y -= 7 * mm

            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, str(label)[:40])
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            s = str(val or '—').replace('\r', ' ').replace('\n', ' ').strip()
            c.drawString(x + label_w, y, s[:int(max(20, value_w / 2))])
            y -= 6.2 * mm

        y -= 4 * mm
        sig_y = max(34 * mm, y)
        self._pdf_draw_signatures(
            c,
            x,
            sig_y,
            width - 2 * margin - 20 * mm,
            left_label='Signature donneur',
            right_label='Signature caissier',
            left_name=None,
            right_name=cashier_name,
        )

        try:
            _draw_authenticity_qr(c, verify_url, tx.receipt_code, x, width, margin, qr_size_mm=26, qr_y_mm=20, text_top_mm=32)
        except Exception:
            pass

        gen_at = timezone.now().strftime('%Y-%m-%d %H:%M')
        c.setFillColorRGB(0.45, 0.50, 0.60)
        c.setFont('Helvetica', 8)
        c.drawRightString(width - margin, 10 * mm, f"Généré le {gen_at}")

        c.showPage()
        c.save()

        pdf = buf.getvalue()
        buf.close()
        return pdf

    def _build_voucher_pdf_bytes(self, tx, request):
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)

        meta = self._pdf_brand_layout(
            c,
            doc_title='BON DE SORTIE',
            doc_subtitle='Finance — Document officiel',
            badge_text='SORTIE',
            badge_rgb=(0.70, 0.12, 0.18),
        )
        width = meta['width']
        margin = meta['margin']
        x = meta['x']
        y = meta['y_start']

        doc_no = tx.document_number or ''

        cashier_name = None
        if tx.cashier:
            cashier_name = (f"{tx.cashier.first_name} {tx.cashier.last_name}").strip() or tx.cashier.username
        cashier_name = cashier_name or '—'

        created_by_name = None
        if tx.created_by:
            created_by_name = (f"{tx.created_by.first_name} {tx.created_by.last_name}").strip() or tx.created_by.username
        created_by_name = created_by_name or '—'

        verify_url = request.build_absolute_uri(f"/api/financial-transactions/verify-document/?code={tx.document_number}")

        def tx_type_label(value):
            try:
                for v, lbl in getattr(FinancialTransaction, 'TRANSACTION_TYPE_CHOICES', []):
                    if v == value:
                        return lbl
            except Exception:
                return value
            return value

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Informations')
        y -= 7 * mm

        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(x, y, 'Numéro')
        c.drawString(x + 55 * mm, y, 'Date')
        y -= 5.5 * mm
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 10)
        c.drawString(x, y, str(doc_no)[:40])
        c.drawString(x + 55 * mm, y, str(tx.date.isoformat())[:10])
        y -= 9 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Destinataire')
        y -= 7 * mm
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 11)
        c.drawString(x, y, str(tx.recipient_name or '—')[:80])
        y -= 7 * mm
        c.setFillColorRGB(0.45, 0.50, 0.60)
        c.setFont('Helvetica', 9)
        c.drawString(x, y, f"Téléphone: {tx.recipient_phone or '—'}")
        y -= 5.5 * mm
        c.drawString(x, y, f"Email: {tx.recipient_email or '—'}")
        y -= 10 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Montant')
        y -= 8 * mm
        c.setFillColorRGB(0.06, 0.18, 0.45)
        c.setFont('Helvetica-Bold', 22)
        c.drawString(x, y, f"{tx.amount} {tx.currency}")
        y -= 12 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 12)
        c.drawString(x, y, 'Détails')
        y -= 7 * mm

        rows = [
            ('Motif', tx_type_label(tx.transaction_type) or '—'),
            ('Mode de paiement', tx.payment_method or '—'),
            ('Référence', tx.reference_number or '—'),
            ('Note', tx.description or '—'),
            ('Fait par', created_by_name),
            ('Approuvé par', cashier_name),
        ]

        label_w = 38 * mm
        value_w = (width - margin - x - 10 * mm) - label_w
        for label, val in rows:
            if y < 32 * mm:
                c.showPage()
                meta = self._pdf_brand_layout(
                    c,
                    doc_title='BON DE SORTIE',
                    doc_subtitle='Finance — Document officiel',
                    badge_text='SORTIE',
                    badge_rgb=(0.70, 0.12, 0.18),
                )
                width = meta['width']
                margin = meta['margin']
                x = meta['x']
                y = meta['y_start']
                c.setFillColorRGB(0.07, 0.10, 0.16)
                c.setFont('Helvetica-Bold', 12)
                c.drawString(x, y, 'Détails')
                y -= 7 * mm

            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, str(label)[:40])
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            s = str(val or '—').replace('\r', ' ').replace('\n', ' ').strip()
            c.drawString(x + label_w, y, s[:int(max(20, value_w / 2))])
            y -= 6.2 * mm

        y -= 4 * mm
        sig_y = max(34 * mm, y)
        self._pdf_draw_signatures(
            c,
            x,
            sig_y,
            width - 2 * margin - 20 * mm,
            left_label='Signature destinataire',
            right_label='Signature caissier',
            left_name=None,
            right_name=cashier_name,
        )

        try:
            _draw_authenticity_qr(c, verify_url, tx.document_number, x, width, margin, qr_size_mm=26, qr_y_mm=20, text_top_mm=32)
        except Exception:
            pass

        gen_at = timezone.now().strftime('%Y-%m-%d %H:%M')
        c.setFillColorRGB(0.45, 0.50, 0.60)
        c.setFont('Helvetica', 8)
        c.drawRightString(width - margin, 10 * mm, f"Généré le {gen_at}")

        c.showPage()
        c.save()

        pdf = buf.getvalue()
        buf.close()
        return pdf

    def _ensure_receipt_pdf(self, tx, request):
        self._ensure_receipt_code(tx)
        if tx.receipt_pdf:
            return tx.receipt_pdf

        pdf = self._build_receipt_pdf_bytes(tx, request)
        filename = f"receipt_{tx.document_number or tx.receipt_code}.pdf"
        tx.receipt_pdf.save(filename, ContentFile(pdf), save=True)
        return tx.receipt_pdf

    def _email_receipt(self, tx, request):
        if not tx.donor_email:
            return False
        if not tx.receipt_pdf:
            return False

        subject = f"Reçu CPD - {tx.receipt_code}"
        donor = tx.donor_name or (tx.member.user.get_full_name() if tx.member and tx.member.user else None) or ''
        body = f"Bonjour {donor},\n\nVeuillez trouver ci-joint votre reçu {tx.receipt_code}.\n\nMerci."
        email = EmailMessage(subject=subject, body=body, from_email=settings.DEFAULT_FROM_EMAIL, to=[tx.donor_email])
        try:
            with open(tx.receipt_pdf.path, 'rb') as f:
                email.attach(f"receipt_{tx.receipt_code}.pdf", f.read(), 'application/pdf')
            email.send(fail_silently=False)
            tx.receipt_sent_at = timezone.now()
            tx.save(update_fields=['receipt_sent_at'])
            return True
        except Exception:
            return False

    def perform_create(self, serializer):
        with transaction.atomic():
            data = dict(serializer.validated_data)
            cashier = data.get('cashier') or self.request.user
            created_by = data.get('created_by') or self.request.user
            tx = serializer.save(cashier=cashier, created_by=created_by)
            self._ensure_document_number(tx)
            if tx.direction == 'in':
                self._ensure_receipt_pdf(tx, self.request)
                if tx.donor_email:
                    self._email_receipt(tx, self.request)

            AuditLogEntry.objects.create(
                actor=self.request.user,
                action='create',
                model='FinancialTransaction',
                object_id=str(tx.pk),
                object_repr=tx.document_number or tx.receipt_code or str(tx.pk),
                ip_address=_client_ip(self.request),
                payload=_safe_payload(self.request.data),
            )

    def perform_update(self, serializer):
        tx = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='FinancialTransaction',
            object_id=str(tx.pk),
            object_repr=tx.document_number or tx.receipt_code or str(tx.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def destroy(self, request, *args, **kwargs):
        tx = self.get_object()
        object_id = str(tx.pk)
        object_repr = tx.document_number or tx.receipt_code or str(tx.pk)
        resp = super().destroy(request, *args, **kwargs)
        if resp.status_code in {200, 202, 204}:
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model='FinancialTransaction',
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
        return resp

    @action(detail=False, methods=['get'], url_path='report')
    def report(self, request):
        period = request.query_params.get('period', 'daily').lower()
        start = request.query_params.get('start')
        end = request.query_params.get('end')

        qs = FinancialTransaction.objects.all()
        if start and end:
            qs = qs.filter(date__range=[start, end])
        elif start:
            qs = qs.filter(date__gte=start)
        elif end:
            qs = qs.filter(date__lte=end)

        if period == 'weekly':
            trunc = TruncWeek('date')
        elif period == 'monthly':
            trunc = TruncMonth('date')
        elif period == 'yearly' or period == 'annual' or period == 'annually':
            trunc = TruncYear('date')
        else:
            trunc = TruncDay('date')

        rows = (
            qs.annotate(p=trunc)
            .values('p', 'currency', 'direction')
            .annotate(total=Sum('amount'))
            .order_by('p', 'currency', 'direction')
        )

        series = {}
        for r in rows:
            key = r['p'].date().isoformat() if hasattr(r['p'], 'date') else str(r['p'])
            currency = r['currency']
            direction = r['direction']
            total = float(r['total'] or 0)

            series.setdefault(key, {})
            series[key].setdefault(currency, {'in': 0.0, 'out': 0.0, 'net': 0.0})
            series[key][currency][direction] = total

        for key, curMap in series.items():
            for cur, agg in curMap.items():
                agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)

        return Response({
            'period': period,
            'start': start,
            'end': end,
            'series': series,
        })

    @action(detail=False, methods=['get'], url_path='summary')
    def summary(self, request):
        rows = (
            self.get_queryset()
            .values('currency', 'direction')
            .annotate(total=Sum('amount'))
            .order_by('currency', 'direction')
        )

        totals = {}
        for r in rows:
            cur = r['currency'] or 'CDF'
            totals.setdefault(cur, {'in': 0.0, 'out': 0.0, 'net': 0.0})
            direction = r['direction']
            totals[cur][direction] = float(r['total'] or 0)

        for cur, agg in totals.items():
            agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)

        return Response({'totals': totals})

    @action(detail=False, methods=['get'], url_path='verify-receipt')
    def verify_receipt(self, request):
        code = (request.query_params.get('code') or '').strip()
        if not code:
            return Response({'detail': 'code requis'}, status=400)

        tx = FinancialTransaction.objects.filter(receipt_code=code, direction='in').select_related('cashier', 'created_by').first()
        if not tx:
            return Response({'detail': 'reçu introuvable'}, status=404)

        return Response({
            'id': tx.id,
            'receipt_code': tx.receipt_code,
            'amount': str(tx.amount),
            'currency': tx.currency,
            'date': tx.date.isoformat(),
            'transaction_type': tx.transaction_type,
            'donor_name': tx.donor_name,
            'donor_email': tx.donor_email,
            'cashier_name': (f"{tx.cashier.first_name} {tx.cashier.last_name}").strip() if tx.cashier else None,
            'created_by_name': (f"{tx.created_by.first_name} {tx.created_by.last_name}").strip() if tx.created_by else None,
            'created_at': tx.created_at.isoformat() if tx.created_at else None,
        })

    @action(detail=False, methods=['get'], url_path='verify-document')
    def verify_document(self, request):
        code = (request.query_params.get('code') or '').strip()
        if not code:
            return Response({'detail': 'code requis'}, status=400)

        tx = FinancialTransaction.objects.filter(document_number=code).select_related('cashier', 'created_by').first()
        if not tx:
            tx = FinancialTransaction.objects.filter(receipt_code=code).select_related('cashier', 'created_by').first()
        if not tx:
            return Response({'detail': 'document introuvable'}, status=404)

        return Response({
            'id': tx.id,
            'document_number': tx.document_number,
            'direction': tx.direction,
            'amount': str(tx.amount),
            'currency': tx.currency,
            'date': tx.date.isoformat(),
            'transaction_type': tx.transaction_type,
            'donor_name': tx.donor_name,
            'donor_email': tx.donor_email,
            'recipient_name': getattr(tx, 'recipient_name', None),
            'recipient_email': getattr(tx, 'recipient_email', None),
            'recipient_phone': getattr(tx, 'recipient_phone', None),
            'cashier_name': (f"{tx.cashier.first_name} {tx.cashier.last_name}").strip() if tx.cashier else None,
            'created_by_name': (f"{tx.created_by.first_name} {tx.created_by.last_name}").strip() if tx.created_by else None,
            'created_at': tx.created_at.isoformat() if tx.created_at else None,
        })

    @action(detail=True, methods=['get'], url_path='receipt')
    def receipt(self, request, pk=None):
        tx = self.get_object()
        if tx.direction != 'in':
            return Response({'detail': 'reçu disponible uniquement pour une entrée'}, status=400)

        self._ensure_receipt_pdf(tx, request)
        if not tx.receipt_pdf:
            return Response({'detail': 'impossible de générer le reçu'}, status=500)

        resp = FileResponse(open(tx.receipt_pdf.path, 'rb'), content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="receipt_{tx.document_number or tx.receipt_code}.pdf"'
        return resp

    @action(detail=True, methods=['get'], url_path='voucher')
    def voucher(self, request, pk=None):
        tx = self.get_object()
        if tx.direction != 'out':
            return Response({'detail': 'bon disponible uniquement pour une sortie'}, status=400)

        with transaction.atomic():
            self._ensure_document_number(tx)
        pdf = self._build_voucher_pdf_bytes(tx, request)

        resp = HttpResponse(pdf, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="voucher_{tx.document_number}.pdf"'
        return resp

    @action(detail=False, methods=['get'], url_path='report-pdf')
    def report_pdf(self, request):
        period = request.query_params.get('period', 'daily').lower()
        start = request.query_params.get('start')
        end = request.query_params.get('end')

        qs = FinancialTransaction.objects.all().order_by('date', 'id')
        if start and end:
            qs = qs.filter(date__range=[start, end])
        elif start:
            qs = qs.filter(date__gte=start)
        elif end:
            qs = qs.filter(date__lte=end)

        if period == 'weekly':
            trunc = TruncWeek('date')
        elif period == 'monthly':
            trunc = TruncMonth('date')
        elif period == 'yearly' or period == 'annual' or period == 'annually':
            trunc = TruncYear('date')
        else:
            trunc = TruncDay('date')

        rows = (
            qs.annotate(p=trunc)
            .values('p', 'currency', 'direction')
            .annotate(total=Sum('amount'))
            .order_by('p', 'currency', 'direction')
        )

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        meta = self._pdf_brand_layout(
            c,
            doc_title='Rapport',
            doc_subtitle='Document officiel — finance',
            badge_text='RAPPORT',
            badge_rgb=(0.34, 0.32, 0.76),
        )

        width = meta['width']
        height = meta['height']
        margin = meta['margin']
        x = meta['x']
        y = meta['y_start']

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 18)
        c.drawString(x, y, 'Bordereau financier')

        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica', 10)
        c.drawRightString(width - margin - 10 * mm, y, f"Période: {start or '—'} → {end or '—'}")
        y -= 10 * mm

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 10)
        period_label = {
            'daily': 'Journalier',
            'weekly': 'Hebdomadaire',
            'monthly': 'Mensuel',
            'yearly': 'Annuel',
            'annual': 'Annuel',
            'annually': 'Annuel',
        }.get(period, period)
        c.drawString(x, y, f"Granularité: {period_label}")
        y -= 10 * mm

        c.setFont('Helvetica-Bold', 10)
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.drawString(x, y, 'Période')
        c.drawString(x + 55 * mm, y, 'Devise')
        c.drawRightString(x + 112 * mm, y, 'Entrées')
        c.drawRightString(x + 138 * mm, y, 'Sorties')
        c.drawRightString(width - margin - 10 * mm, y, 'Solde')

        y -= 6 * mm
        c.setFont('Helvetica', 10)
        c.setFillColorRGB(0.20, 0.24, 0.30)
        series = {}
        for r in rows:
            key = r['p'].date().isoformat() if hasattr(r['p'], 'date') else str(r['p'])
            cur = r['currency']
            direction = r['direction']
            total = float(r['total'] or 0)
            series.setdefault(key, {})
            series[key].setdefault(cur, {'in': 0.0, 'out': 0.0, 'net': 0.0})
            series[key][cur][direction] = total

        for key, curMap in series.items():
            for cur, agg in curMap.items():
                agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)

        for p, curMap in series.items():
            for cur, agg in curMap.items():
                if y < 20 * mm:
                    c.showPage()
                    meta_p = self._pdf_brand_layout(
                        c,
                        doc_title='Rapport',
                        doc_subtitle='Document officiel — finance',
                        badge_text='RAPPORT',
                        badge_rgb=(0.34, 0.32, 0.76),
                    )
                    width = meta_p['width']
                    height = meta_p['height']
                    margin = meta_p['margin']
                    x = meta_p['x']
                    y = meta_p['y_start']
                    c.setFont('Helvetica', 10)
                    c.setFillColorRGB(0.20, 0.24, 0.30)

                c.drawString(x, y, str(p))
                c.drawString(x + 55 * mm, y, str(cur))
                c.drawRightString(x + 112 * mm, y, f"{agg.get('in', 0):.2f}")
                c.drawRightString(x + 138 * mm, y, f"{agg.get('out', 0):.2f}")
                c.drawRightString(width - margin - 10 * mm, y, f"{agg.get('net', 0):.2f}")
                y -= 6 * mm

        c.showPage()
        c.save()
        pdf = buf.getvalue()
        buf.close()

        resp = HttpResponse(pdf, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="report_{period}.pdf"'
        return resp


class AnnouncementViewSet(viewsets.ModelViewSet):
    queryset = Announcement.objects.all().order_by('-published_date')
    serializer_class = AnnouncementSerializer
    permission_classes = [PublicReadAdminWrite]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_permissions(self):
        if getattr(self, 'action', None) in {'like', 'comments', 'comment_like'}:
            return [IsAuthenticated()]
        return super().get_permissions()

    def perform_create(self, serializer):
        serializer.save(author=self.request.user)

    @action(detail=True, methods=['post'], url_path='like')
    def like(self, request, pk=None):
        ann = self.get_object()
        existing = AnnouncementLike.objects.filter(announcement=ann, user=request.user).first()
        if existing:
            existing.delete()
            liked = False
        else:
            AnnouncementLike.objects.create(announcement=ann, user=request.user)
            liked = True
        return Response({'liked': liked, 'like_count': AnnouncementLike.objects.filter(announcement=ann).count()})

    @action(detail=True, methods=['get', 'post'], url_path='comments')
    def comments(self, request, pk=None):
        ann = self.get_object()
        if request.method == 'GET':
            qs = AnnouncementComment.objects.filter(announcement=ann).order_by('-created_at')
            return Response(AnnouncementCommentSerializer(qs, many=True, context={'request': request}).data)

        serializer = AnnouncementCommentSerializer(data=request.data, context={'request': request})
        serializer.is_valid(raise_exception=True)
        serializer.save(announcement=ann, author=request.user)
        return Response(serializer.data, status=201)

    @action(detail=True, methods=['post'], url_path=r'comments/(?P<comment_id>\d+)/like')
    def comment_like(self, request, pk=None, comment_id=None):
        ann = self.get_object()
        comment = AnnouncementComment.objects.filter(id=comment_id, announcement=ann).first()
        if not comment:
            return Response({'detail': 'commentaire introuvable'}, status=404)

        existing = AnnouncementCommentLike.objects.filter(comment=comment, user=request.user).first()
        if existing:
            existing.delete()
            liked = False
        else:
            AnnouncementCommentLike.objects.create(comment=comment, user=request.user)
            liked = True

        return Response({'liked': liked, 'like_count': AnnouncementCommentLike.objects.filter(comment=comment).count()})


class AnnouncementDeckViewSet(viewsets.ModelViewSet):
    queryset = AnnouncementDeck.objects.select_related('event', 'created_by').prefetch_related('items').all().order_by('-created_at', '-id')
    serializer_class = AnnouncementDeckSerializer
    permission_classes = [IsAdminOrSuperAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)

    def perform_update(self, serializer):
        serializer.save()

    @action(detail=True, methods=['get', 'post'], url_path='items')
    def items(self, request, pk=None):
        deck = self.get_object()

        if request.method == 'GET':
            qs = getattr(deck, 'items', None).all() if getattr(deck, 'items', None) is not None else AnnouncementDeckItem.objects.filter(deck=deck)
            return Response(AnnouncementDeckItemSerializer(qs, many=True).data)

        serializer = AnnouncementDeckItemSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        serializer.save(deck=deck)
        return Response(serializer.data, status=201)

    @action(detail=True, methods=['post'], url_path='set-items')
    def set_items(self, request, pk=None):
        deck = self.get_object()
        raw = request.data.get('items')
        if not isinstance(raw, list):
            return Response({'detail': 'items doit être une liste.'}, status=400)

        cleaned = []
        for idx, it in enumerate(raw, start=1):
            if not isinstance(it, dict):
                continue
            text = (it.get('text') or '').strip()
            if not text:
                continue
            try:
                order = int(it.get('order') or idx)
            except Exception:
                order = idx
            cleaned.append({'order': max(1, order), 'text': text})

        with transaction.atomic():
            AnnouncementDeckItem.objects.filter(deck=deck).delete()
            for it in cleaned:
                AnnouncementDeckItem.objects.create(deck=deck, order=it['order'], text=it['text'])

        qs = AnnouncementDeckItem.objects.filter(deck=deck).order_by('order', 'id')
        return Response(AnnouncementDeckItemSerializer(qs, many=True).data)

    @action(detail=True, methods=['post'], url_path='generate')
    def generate(self, request, pk=None):
        deck = self.get_object()

        if Presentation is None:
            return Response({'detail': 'Génération PPTX indisponible: installe python-pptx sur le backend.'}, status=503)

        items = list(getattr(deck, 'items', None).all()) if getattr(deck, 'items', None) is not None else list(AnnouncementDeckItem.objects.filter(deck=deck))
        header = (getattr(deck, 'header_text', None) or '').strip()
        theme = (getattr(deck, 'theme_text', None) or '').strip()

        ev = getattr(deck, 'event', None)
        ev_title = getattr(ev, 'title', None) if ev else None
        ev_date = getattr(ev, 'date', None) if ev else None
        ev_time = getattr(ev, 'time', None) if ev else None
        ev_location = getattr(ev, 'location', None) if ev else None

        subtitle_parts = []
        if ev_date:
            subtitle_parts.append(str(ev_date))
        if ev_time:
            subtitle_parts.append(str(ev_time)[:5])
        if ev_location:
            subtitle_parts.append(str(ev_location))
        if theme:
            subtitle_parts.append(str(theme))
        subtitle = ' • '.join([p for p in subtitle_parts if str(p).strip()])

        prs = Presentation()
        title = (getattr(deck, 'title', None) or ev_title or 'Annonces')
        _pptx_add_title(prs, title, subtitle=subtitle)

        if header:
            _pptx_add_announcement_slide(prs, header='En-tête', text=header, number=None)

        for idx, it in enumerate(items, start=1):
            _pptx_add_announcement_slide(prs, header=title, text=getattr(it, 'text', None), number=idx)

        buf = io.BytesIO()
        prs.save(buf)
        pptx = buf.getvalue()
        buf.close()

        stamp = timezone.now().strftime('%Y%m%d_%H%M%S')
        safe_title = re.sub(r'[^a-zA-Z0-9_-]+', '_', str(title))[:60].strip('_') or 'annonces'
        filename = f"deck_{safe_title}_{stamp}.pptx"

        deck.pptx_file.save(filename, ContentFile(pptx), save=False)
        deck.generated_at = timezone.now()
        deck.save(update_fields=['pptx_file', 'generated_at', 'updated_at'])

        resp = HttpResponse(pptx, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp

    @action(detail=True, methods=['get'], url_path='download')
    def download(self, request, pk=None):
        deck = self.get_object()
        f = getattr(deck, 'pptx_file', None)
        if not f:
            return Response({'detail': 'Fichier PPTX indisponible. Génère le deck.'}, status=404)
        try:
            return FileResponse(f.open('rb'), as_attachment=True, filename=getattr(f, 'name', None) or 'annonces.pptx')
        except Exception:
            return Response({'detail': 'Impossible de télécharger le fichier.'}, status=500)


class ReportViewSet(viewsets.ViewSet):
    permission_classes = [IsAdminOrSuperAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_permissions(self):
        if getattr(self, 'action', None) in {'verify'}:
            return [AllowAny()]
        return super().get_permissions()

    def list(self, request):
        return Response({
            'available_periods': ['daily', 'weekly', 'monthly', 'annual'],
            'available_sections': ['programmes', 'pointage', 'members', 'finances', 'diaconat', 'evangelisation'],
            'default_sections': ['programmes', 'pointage', 'members', 'finances', 'diaconat', 'evangelisation'],
            'endpoints': {
                'compiled': '/api/reports/compiled/',
                'compiled_pdf': '/api/reports/compiled-pdf/',
                'verify': '/api/reports/verify/?code=...',
            },
        })

    @action(detail=False, methods=['get'], url_path='verify')
    def verify(self, request):
        code = (request.query_params.get('code') or '').strip()
        if not code:
            return Response({'valid': False, 'detail': 'code requis'}, status=400)

        cert = ReportCertificate.objects.filter(code=code).first()
        if not cert:
            return Response({'valid': False, 'detail': 'certificat introuvable'}, status=404)

        return Response({
            'valid': True,
            'code': cert.code,
            'report_type': cert.report_type,
            'payload': cert.payload,
            'pdf_sha256': cert.pdf_sha256,
            'created_at': cert.created_at.isoformat() if getattr(cert, 'created_at', None) else None,
        })

    def _parse_sections(self, request):
        raw = []
        try:
            raw = list(request.query_params.getlist('sections') or [])
        except Exception:
            raw = []

        if not raw:
            s = (request.query_params.get('sections') or '').strip()
            if s:
                raw = [s]

        parts = []
        for item in raw:
            if not item:
                continue
            if isinstance(item, str) and ',' in item:
                parts.extend([x.strip() for x in item.split(',') if x.strip()])
            else:
                parts.append(str(item).strip())

        allowed = {'programmes', 'members', 'finances', 'logistics', 'diaconat', 'pointage', 'evangelisation'}
        out = []
        for p in parts:
            key = p.lower()
            if key in allowed and key not in out:
                out.append(key)
        return out

    def _diaconat_section(self, start_d, end_d):
        qs = LogisticsItem.objects.all().order_by('-updated_at', '-id')
        in_range = qs.filter(Q(acquired_date__range=[start_d, end_d]) | Q(created_at__date__range=[start_d, end_d]))

        by_category = list(
            qs.values('category')
            .annotate(items=Count('id'), quantity=Sum('quantity'))
            .order_by('-items', 'category')
        )
        by_condition = list(
            qs.values('condition')
            .annotate(items=Count('id'), quantity=Sum('quantity'))
            .order_by('-items', 'condition')
        )

        purchase_total = qs.aggregate(total=Sum('purchase_price'))
        purchase_in_range = in_range.aggregate(total=Sum('purchase_price'))

        return {
            'items_count': qs.count(),
            'active_items_count': qs.filter(is_active=True).count(),
            'quantity_total': float(qs.aggregate(q=Sum('quantity')).get('q') or 0),
            'by_category': by_category,
            'by_condition': by_condition,
            'purchase_total': float(purchase_total.get('total') or 0),
            'purchase_total_in_period': float(purchase_in_range.get('total') or 0),
            'created_or_acquired_in_period_count': in_range.count(),
        }

    def _evangelisation_section(self, start_d, end_d):
        bap_qs = BaptismEvent.objects.select_related('event').filter(event__date__range=[start_d, end_d])
        ev_qs = EvangelismActivity.objects.filter(date__range=[start_d, end_d])
        tr_qs = TrainingEvent.objects.filter(date__range=[start_d, end_d])
        cand_qs = BaptismCandidate.objects.filter(baptism_event__event__date__range=[start_d, end_d])

        return {
            'baptisms_count': bap_qs.count(),
            'candidates_count': cand_qs.count(),
            'evangelism_activities_count': ev_qs.count(),
            'training_events_count': tr_qs.count(),
        }

    def _infer_range(self, period, start, end):
        today = timezone.localdate()

        def parse_date(s):
            if not s:
                return None
            try:
                return datetime.date.fromisoformat(str(s)[:10])
            except Exception:
                return None

        start_d = parse_date(start)
        end_d = parse_date(end)

        if start_d and not end_d:
            end_d = start_d
        if end_d and not start_d:
            start_d = end_d

        if start_d and end_d:
            return start_d, end_d

        p = (period or 'daily').lower()
        if p in {'yearly', 'annual', 'annually'}:
            start_d = datetime.date(today.year, 1, 1)
            end_d = datetime.date(today.year, 12, 31)
        elif p == 'monthly':
            start_d = today.replace(day=1)
            next_month = (start_d + datetime.timedelta(days=32)).replace(day=1)
            end_d = next_month - datetime.timedelta(days=1)
        elif p == 'weekly':
            start_d = today - datetime.timedelta(days=today.weekday())
            end_d = start_d + datetime.timedelta(days=6)
        else:
            start_d = today
            end_d = today

        return start_d, end_d

    def _period_trunc(self, period):
        p = (period or 'daily').lower()
        if p == 'weekly':
            return TruncWeek
        if p == 'monthly':
            return TruncMonth
        if p in {'yearly', 'annual', 'annually'}:
            return TruncYear
        return TruncDay

    def _programmes_section(self, start_d, end_d):
        qs = Event.objects.filter(date__range=[start_d, end_d]).order_by('date', 'time', 'id')

        by_type = list(
            qs.values('event_type')
            .annotate(count=Count('id'))
            .order_by('-count', 'event_type')
        )

        items = []
        for ev in qs[:250]:
            items.append({
                'id': ev.id,
                'title': ev.title,
                'date': ev.date.isoformat() if getattr(ev, 'date', None) else None,
                'time': ev.time.strftime('%H:%M') if getattr(ev, 'time', None) else None,
                'event_type': ev.event_type,
                'duration_type': ev.duration_type,
                'department_name': ev.department.name if getattr(ev, 'department', None) else None,
                'is_published': bool(getattr(ev, 'is_published', False)),
            })

        return {
            'count': qs.count(),
            'published_count': qs.filter(is_published=True).count(),
            'department_events_count': qs.filter(department__isnull=False).count(),
            'by_type': by_type,
            'events': items,
        }

    def _members_section(self, start_d, end_d):
        qs = Member.objects.select_related('user', 'department', 'ministry').all()
        new_qs = qs.filter(created_at__date__range=[start_d, end_d]).order_by('-created_at', '-id')

        by_gender = list(
            qs.values('gender')
            .annotate(count=Count('id'))
            .order_by('-count', 'gender')
        )

        by_department = list(
            qs.values('department__name')
            .annotate(count=Count('id'))
            .order_by('-count', 'department__name')
        )

        recent = []
        for m in new_qs[:30]:
            u = getattr(m, 'user', None)
            recent.append({
                'id': m.id,
                'member_number': m.member_number,
                'full_name': (u.get_full_name() if u else '').strip() or (u.username if u else ''),
                'phone': getattr(u, 'phone', None) if u else None,
                'gender': m.gender,
                'department_name': m.department.name if getattr(m, 'department', None) else None,
                'ministry_name': m.ministry.name if getattr(m, 'ministry', None) else None,
                'created_at': m.created_at.isoformat() if getattr(m, 'created_at', None) else None,
            })

        return {
            'total': qs.count(),
            'active': qs.filter(is_active=True).count(),
            'inactive': qs.filter(is_active=False).count(),
            'new_count': new_qs.count(),
            'by_gender': by_gender,
            'by_department': by_department,
            'recent_new_members': recent,
        }

    def _finances_section(self, start_d, end_d):
        qs = FinancialTransaction.objects.filter(date__range=[start_d, end_d]).order_by('date', 'id')

        rows = list(
            qs.values('currency', 'direction')
            .annotate(total=Sum('amount'))
            .order_by('currency', 'direction')
        )
        totals = {}
        for r in rows:
            cur = r['currency'] or 'CDF'
            totals.setdefault(cur, {'in': 0.0, 'out': 0.0, 'net': 0.0})
            direction = r['direction']
            totals[cur][direction] = float(r['total'] or 0)
        for cur, agg in totals.items():
            agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)

        by_type_rows = list(
            qs.values('currency', 'transaction_type', 'direction')
            .annotate(total=Sum('amount'))
            .order_by('currency', 'transaction_type', 'direction')
        )
        by_type = {}
        for r in by_type_rows:
            cur = r['currency'] or 'CDF'
            tx_type = r['transaction_type'] or '—'
            direction = r['direction']
            by_type.setdefault(cur, {})
            by_type[cur].setdefault(tx_type, {'in': 0.0, 'out': 0.0, 'net': 0.0})
            by_type[cur][tx_type][direction] = float(r['total'] or 0)
        for cur, mp in by_type.items():
            for tx_type, agg in mp.items():
                agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)

        breakdown_rows = list(
            qs.values('currency', 'transaction_type', 'description', 'direction')
            .annotate(total=Sum('amount'))
            .order_by('currency', 'transaction_type', 'direction', 'description')
        )
        breakdown = []
        for r in breakdown_rows[:400]:
            breakdown.append({
                'currency': r['currency'] or 'CDF',
                'transaction_type': r['transaction_type'] or '—',
                'description': (r.get('description') or '').strip() or '—',
                'direction': r['direction'] or 'in',
                'total': float(r['total'] or 0),
            })

        activity_qs = qs.filter(event__isnull=False)
        non_activity_qs = qs.filter(event__isnull=True)

        def dir_totals(tx_qs):
            res = list(tx_qs.values('currency', 'direction').annotate(total=Sum('amount')).order_by('currency', 'direction'))
            out = {}
            for r in res:
                cur = r['currency'] or 'CDF'
                out.setdefault(cur, {'in': 0.0, 'out': 0.0, 'net': 0.0})
                out[cur][r['direction']] = float(r['total'] or 0)
            for cur, agg in out.items():
                agg['net'] = float(agg.get('in', 0) or 0) - float(agg.get('out', 0) or 0)
            return out

        return {
            'transaction_count': qs.count(),
            'totals': totals,
            'by_type': by_type,
            'breakdown': breakdown,
            'activity_related': {
                'transaction_count': activity_qs.count(),
                'totals': dir_totals(activity_qs),
            },
            'non_activity_related': {
                'transaction_count': non_activity_qs.count(),
                'totals': dir_totals(non_activity_qs),
            },
        }

    def _logistics_section(self, start_d, end_d):
        qs = LogisticsItem.objects.all().order_by('-updated_at', '-id')
        in_range = qs.filter(Q(acquired_date__range=[start_d, end_d]) | Q(created_at__date__range=[start_d, end_d]))

        by_category = list(
            qs.values('category')
            .annotate(items=Count('id'), quantity=Sum('quantity'))
            .order_by('-items', 'category')
        )
        by_condition = list(
            qs.values('condition')
            .annotate(items=Count('id'), quantity=Sum('quantity'))
            .order_by('-items', 'condition')
        )

        purchase_total = qs.aggregate(total=Sum('purchase_price'))
        purchase_in_range = in_range.aggregate(total=Sum('purchase_price'))

        return {
            'items_count': qs.count(),
            'active_items_count': qs.filter(is_active=True).count(),
            'quantity_total': float(qs.aggregate(q=Sum('quantity')).get('q') or 0),
            'by_category': by_category,
            'by_condition': by_condition,
            'purchase_total': float(purchase_total.get('total') or 0),
            'purchase_total_in_period': float(purchase_in_range.get('total') or 0),
            'created_or_acquired_in_period_count': in_range.count(),
        }

    def _pointage_section(self, start_d, end_d, period):
        ev_qs = Event.objects.filter(date__range=[start_d, end_d])
        agg_qs = EventAttendanceAggregate.objects.filter(event__date__range=[start_d, end_d]).select_related('event')

        totals = agg_qs.aggregate(
            male_adults=Sum('male_adults'),
            female_adults=Sum('female_adults'),
            male_children=Sum('male_children'),
            female_children=Sum('female_children'),
        )
        totals = {k: int(v or 0) for k, v in totals.items()}
        totals['total'] = int(totals.get('male_adults', 0) + totals.get('female_adults', 0) + totals.get('male_children', 0) + totals.get('female_children', 0))

        attendance_qs = Attendance.objects.filter(event__date__range=[start_d, end_d])
        attended_count = attendance_qs.filter(attended=True).count()
        total_rows = attendance_qs.count()
        rate = (float(attended_count) / float(total_rows)) if total_rows else 0.0

        by_dept = list(
            attendance_qs.filter(attended=True, member__department__isnull=False)
            .values('member__department__name')
            .annotate(count=Count('id'))
            .order_by('-count', 'member__department__name')
        )

        trunc = self._period_trunc(period)
        series_rows = list(
            agg_qs.annotate(p=trunc('event__date'))
            .values('p')
            .annotate(
                male_adults=Sum('male_adults'),
                female_adults=Sum('female_adults'),
                male_children=Sum('male_children'),
                female_children=Sum('female_children'),
            )
            .order_by('p')
        )
        series = []
        for r in series_rows:
            p = r.get('p')
            key = p.date().isoformat() if hasattr(p, 'date') else str(p)
            t = int((r.get('male_adults') or 0) + (r.get('female_adults') or 0) + (r.get('male_children') or 0) + (r.get('female_children') or 0))
            series.append({'period': key, 'total': t})

        by_event_type_rows = list(
            agg_qs.values('event__event_type')
            .annotate(
                male_adults=Sum('male_adults'),
                female_adults=Sum('female_adults'),
                male_children=Sum('male_children'),
                female_children=Sum('female_children'),
            )
            .order_by('event__event_type')
        )
        by_event_type = []
        for r in by_event_type_rows:
            t = int((r.get('male_adults') or 0) + (r.get('female_adults') or 0) + (r.get('male_children') or 0) + (r.get('female_children') or 0))
            by_event_type.append({'event_type': r.get('event__event_type') or '—', 'total': t})

        return {
            'events_count': ev_qs.count(),
            'attendance_rows_count': total_rows,
            'members_present_count': attended_count,
            'attendance_rate': rate,
            'anonymous_totals': totals,
            'by_department': by_dept,
            'series': series,
            'by_event_type': by_event_type,
        }

    def _compiled_report_data(self, request):
        period = (request.query_params.get('period') or 'daily').lower()
        start = request.query_params.get('start')
        end = request.query_params.get('end')
        sections = self._parse_sections(request)
        if not sections:
            sections = ['programmes', 'pointage', 'members', 'finances', 'diaconat', 'evangelisation']

        # Backward compatibility
        if 'logistics' in sections and 'diaconat' not in sections:
            sections.append('diaconat')

        start_d, end_d = self._infer_range(period, start, end)

        payload = {
            'period': period,
            'start': start_d.isoformat(),
            'end': end_d.isoformat(),
            'generated_at': timezone.now().isoformat(),
            'sections': sections,
            'data': {},
        }

        if 'programmes' in sections:
            payload['data']['programmes'] = self._programmes_section(start_d, end_d)
        if 'pointage' in sections:
            payload['data']['pointage'] = self._pointage_section(start_d, end_d, period)
        if 'members' in sections:
            payload['data']['members'] = self._members_section(start_d, end_d)
        if 'finances' in sections:
            payload['data']['finances'] = self._finances_section(start_d, end_d)
        if 'diaconat' in sections or 'logistics' in sections:
            payload['data']['diaconat'] = self._diaconat_section(start_d, end_d)
        if 'logistics' in sections:
            payload['data']['logistics'] = payload['data'].get('diaconat')
        if 'evangelisation' in sections:
            payload['data']['evangelisation'] = self._evangelisation_section(start_d, end_d)

        return payload

    def _compiled_report_pdf(self, request, compiled):
        cert = None
        verify_url = None
        try:
            payload = {
                'period': compiled.get('period'),
                'start': compiled.get('start'),
                'end': compiled.get('end'),
                'sections': compiled.get('sections'),
                'generated_at': timezone.now().isoformat(),
            }
            cert = _create_report_certificate('compiled', payload, getattr(request, 'user', None))
            verify_url = request.build_absolute_uri(f"/api/reports/verify/?code={cert.code}")
        except Exception:
            cert = None
            verify_url = None

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        meta = FinancialTransactionViewSet()._pdf_brand_layout(
            c,
            doc_title='Rapport',
            doc_subtitle='Document officiel — rapport compilé',
            badge_text='GLOBAL',
            badge_rgb=(0.20, 0.35, 0.75),
        )

        width = meta['width']
        margin = meta['margin']
        x = meta['x']
        y = meta['y_start']

        def new_page():
            nonlocal width, margin, x, y
            c.showPage()
            meta2 = FinancialTransactionViewSet()._pdf_brand_layout(
                c,
                doc_title='Rapport',
                doc_subtitle='Document officiel — rapport compilé',
                badge_text='GLOBAL',
                badge_rgb=(0.20, 0.35, 0.75),
            )
            width = meta2['width']
            margin = meta2['margin']
            x = meta2['x']
            y = meta2['y_start']

        def ensure(min_y=25 * mm):
            nonlocal y
            if y < min_y:
                new_page()

        c.setFillColorRGB(0.07, 0.10, 0.16)
        c.setFont('Helvetica-Bold', 16)
        c.drawString(x, y, 'Rapport global (compilé)')
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica', 10)
        c.drawRightString(width - margin - 10 * mm, y, f"Période: {compiled.get('start')} → {compiled.get('end')}")
        y -= 10 * mm

        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(x, y, 'Sections')
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 10)
        c.drawString(x + 18 * mm, y, ', '.join(compiled.get('sections') or []))
        y -= 10 * mm

        data = compiled.get('data') or {}

        if 'programmes' in data:
            ensure()
            p = data['programmes'] or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 12)
            c.drawString(x, y, 'Programmes & activités')
            y -= 7 * mm
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            c.drawString(x, y, f"Nombre d'activités: {p.get('count', 0)}")
            c.drawRightString(width - margin - 10 * mm, y, f"Publiées: {p.get('published_count', 0)}")
            y -= 6 * mm
            c.drawString(x, y, f"Activités de département: {p.get('department_events_count', 0)}")
            y -= 8 * mm

        if 'pointage' in data:
            ensure()
            pt = data['pointage'] or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 12)
            c.drawString(x, y, 'Pointage & statistiques')
            y -= 7 * mm
            anon = pt.get('anonymous_totals') or {}
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            c.drawString(x, y, f"Total (anonyme): {anon.get('total', 0)}")
            c.drawRightString(width - margin - 10 * mm, y, f"Activités: {pt.get('events_count', 0)}")
            y -= 6 * mm
            c.drawString(x, y, f"Présence membres: {pt.get('members_present_count', 0)} (taux: {pt.get('attendance_rate', 0.0) * 100:.1f}%)")
            y -= 8 * mm

        if 'members' in data:
            ensure()
            m = data['members'] or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 12)
            c.drawString(x, y, 'Membres')
            y -= 7 * mm
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            c.drawString(x, y, f"Total: {m.get('total', 0)}")
            c.drawRightString(width - margin - 10 * mm, y, f"Nouveaux (période): {m.get('new_count', 0)}")
            y -= 6 * mm
            c.drawString(x, y, f"Actifs: {m.get('active', 0)}  •  Inactifs: {m.get('inactive', 0)}")
            y -= 8 * mm

        if 'finances' in data:
            ensure()
            f = data['finances'] or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 12)
            c.drawString(x, y, 'Finances')
            y -= 7 * mm
            c.setFillColorRGB(0.40, 0.45, 0.55)
            c.setFont('Helvetica-Bold', 9)
            c.drawString(x, y, 'Devise')
            c.drawRightString(x + 95 * mm, y, 'Entrées')
            c.drawRightString(x + 125 * mm, y, 'Sorties')
            c.drawRightString(width - margin - 10 * mm, y, 'Solde')
            y -= 6 * mm
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            for cur, agg in (f.get('totals') or {}).items():
                ensure()
                c.drawString(x, y, str(cur))
                c.drawRightString(x + 95 * mm, y, f"{float(agg.get('in', 0) or 0):.2f}")
                c.drawRightString(x + 125 * mm, y, f"{float(agg.get('out', 0) or 0):.2f}")
                c.drawRightString(width - margin - 10 * mm, y, f"{float(agg.get('net', 0) or 0):.2f}")
                y -= 6 * mm
            y -= 4 * mm

        if 'logistics' in data:
            ensure()
            l = data['logistics'] or {}
            c.setFillColorRGB(0.07, 0.10, 0.16)
            c.setFont('Helvetica-Bold', 12)
            c.drawString(x, y, 'Logistique')
            y -= 7 * mm
            c.setFillColorRGB(0.20, 0.24, 0.30)
            c.setFont('Helvetica', 10)
            c.drawString(x, y, f"Articles: {l.get('items_count', 0)}  •  Actifs: {l.get('active_items_count', 0)}")
            y -= 6 * mm
            c.drawString(x, y, f"Quantité totale: {float(l.get('quantity_total', 0) or 0):.0f}  •  Valeur d'achat: {float(l.get('purchase_total', 0) or 0):.2f}")
            y -= 8 * mm

        if y < 90 * mm:
            new_page()

        sig_x = x
        sig_y = 56 * mm
        sig_w = 85 * mm
        c.setFillColorRGB(0.40, 0.45, 0.55)
        c.setFont('Helvetica-Bold', 9)
        c.drawString(sig_x, sig_y + 6 * mm, 'Signature')
        c.setFillColorRGB(0.20, 0.24, 0.30)
        c.setFont('Helvetica', 9)
        c.drawString(sig_x, sig_y + 2.5 * mm, 'Administrateur')
        c.line(sig_x, sig_y, sig_x + sig_w, sig_y)

        if cert and verify_url:
            try:
                _draw_authenticity_qr(c, verify_url, cert.code, x, width, margin, qr_size_mm=26)
            except Exception:
                pass

        c.setFillColorRGB(0.45, 0.50, 0.60)
        c.setFont('Helvetica', 8)
        gen_at = timezone.now().strftime('%Y-%m-%d %H:%M')
        c.drawRightString(width - margin, 10 * mm, f"Généré le {gen_at}")

        c.showPage()
        c.save()
        pdf = buf.getvalue()
        buf.close()

        if cert:
            try:
                cert.pdf_sha256 = hashlib.sha256(pdf).hexdigest()
                cert.save(update_fields=['pdf_sha256', 'updated_at'])
            except Exception:
                pass

        return pdf

    @action(detail=False, methods=['get'], url_path='compiled')
    def compiled(self, request):
        return Response(self._compiled_report_data(request))

    @action(detail=False, methods=['get'], url_path='compiled-pdf')
    def compiled_pdf(self, request):
        compiled = self._compiled_report_data(request)
        pdf = self._compiled_report_pdf(request, compiled)
        stamp = timezone.now().strftime('%Y%m%d_%H%M%S')
        filename = f"rapport_global_{stamp}.pdf"
        resp = HttpResponse(pdf, content_type='application/pdf')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp


class DocumentViewSet(viewsets.ModelViewSet):
    queryset = Document.objects.all().order_by('-uploaded_at')
    serializer_class = DocumentSerializer
    permission_classes = [IsAdminOrSuperAdminOrReadOnly]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_queryset(self):
        qs = super().get_queryset()
        q = (self.request.query_params.get('q') or '').strip()
        doc_type = (self.request.query_params.get('document_type') or '').strip()
        if doc_type:
            qs = qs.filter(document_type=doc_type)
        if q:
            qs = qs.filter(Q(title__icontains=q))
        return qs

    def perform_create(self, serializer):
        obj = serializer.save(uploaded_by=self.request.user)
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='Document',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'title', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def perform_update(self, serializer):
        obj = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='Document',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'title', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def destroy(self, request, *args, **kwargs):
        obj = self.get_object()
        object_id = str(obj.pk)
        object_repr = getattr(obj, 'title', None) or str(obj.pk)
        resp = super().destroy(request, *args, **kwargs)
        if resp.status_code in {200, 202, 204}:
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model='Document',
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
        return resp


class LogisticsItemViewSet(viewsets.ModelViewSet):
    queryset = LogisticsItem.objects.all().order_by('-id')
    serializer_class = LogisticsItemSerializer
    permission_classes = [IsLogisticsHeadOrAdmin]

    def create(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().create(request, *args, **kwargs)
        ar = _create_approval_request(request, model='LogisticsItem', action='create', payload=_safe_payload(request.data))
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def update(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            return super().update(request, *args, **kwargs)
        obj = self.get_object()
        ar = _create_approval_request(
            request,
            model='LogisticsItem',
            action='update',
            payload=_safe_payload(request.data),
            target_object_id=getattr(obj, 'id', None),
            object_repr=getattr(obj, 'name', None) or str(getattr(obj, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def partial_update(self, request, *args, **kwargs):
        return self.update(request, *args, **kwargs)

    def destroy(self, request, *args, **kwargs):
        if _is_admin_user(request.user):
            obj = self.get_object()
            object_id = str(obj.pk)
            object_repr = getattr(obj, 'name', None) or str(obj.pk)
            if getattr(obj, 'is_active', True):
                obj.is_active = False
                obj.save(update_fields=['is_active', 'updated_at'])
            AuditLogEntry.objects.create(
                actor=request.user,
                action='delete',
                model='LogisticsItem',
                object_id=object_id,
                object_repr=object_repr,
                ip_address=_client_ip(request),
                payload=None,
            )
            return Response(status=204)
        obj = self.get_object()
        ar = _create_approval_request(
            request,
            model='LogisticsItem',
            action='delete',
            payload=None,
            target_object_id=getattr(obj, 'id', None),
            object_repr=getattr(obj, 'name', None) or str(getattr(obj, 'id', '') or ''),
        )
        return Response({'detail': 'Action soumise à approbation.', 'approval_request_id': ar.id}, status=202)

    def get_queryset(self):
        qs = super().get_queryset()
        qs = qs.filter(is_active=True)
        q = (self.request.query_params.get('q') or '').strip()
        if q:
            qs = qs.filter(
                Q(name__icontains=q)
                | Q(category__icontains=q)
                | Q(asset_tag__icontains=q)
                | Q(location__icontains=q)
                | Q(supplier__icontains=q)
            )
        return qs

    def perform_create(self, serializer):
        obj = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='create',
            model='LogisticsItem',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'name', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )

    def perform_update(self, serializer):
        obj = serializer.save()
        AuditLogEntry.objects.create(
            actor=self.request.user,
            action='update',
            model='LogisticsItem',
            object_id=str(obj.pk),
            object_repr=getattr(obj, 'name', None) or str(obj.pk),
            ip_address=_client_ip(self.request),
            payload=_safe_payload(self.request.data),
        )


class AuditLogEntryViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = AuditLogEntry.objects.select_related('actor').all().order_by('-id')
    serializer_class = AuditLogEntrySerializer
    permission_classes = [IsAdminOrSuperAdmin]

    def get_queryset(self):
        qs = super().get_queryset()
        model = (self.request.query_params.get('model') or '').strip()
        action = (self.request.query_params.get('action') or '').strip()
        q = (self.request.query_params.get('q') or '').strip()
        if model:
            qs = qs.filter(model__iexact=model)
        if action:
            qs = qs.filter(action__iexact=action)
        if q:
            qs = qs.filter(Q(object_id__icontains=q) | Q(object_repr__icontains=q) | Q(actor__username__icontains=q))
        return qs


class NotificationViewSet(viewsets.ModelViewSet):
    queryset = Notification.objects.all().order_by('-id')
    serializer_class = NotificationSerializer
    permission_classes = [IsAuthenticated]

    def get_queryset(self):
        if self.request.user.is_superuser:
            return Notification.objects.all().order_by('-id')
        return Notification.objects.filter(recipient=self.request.user).order_by('-id')


class ApprovalRequestViewSet(viewsets.ModelViewSet):
    queryset = ApprovalRequest.objects.select_related('requested_by', 'decided_by').all().order_by('-created_at', '-id')
    serializer_class = ApprovalRequestSerializer
    permission_classes = [IsAuthenticated]

    def get_queryset(self):
        qs = super().get_queryset()
        u = getattr(self.request, 'user', None)
        if not u or not getattr(u, 'is_authenticated', False):
            return qs.none()
        if _is_admin_user(u):
            status = (self.request.query_params.get('status') or '').strip().lower()
            if status in {'pending', 'approved', 'rejected'}:
                qs = qs.filter(status=status)
            return qs
        return qs.filter(requested_by=u)

    def create(self, request, *args, **kwargs):
        raise PermissionDenied('Création manuelle non autorisée.')

    @action(detail=True, methods=['post'], url_path='approve', permission_classes=[IsAdminOrSuperAdmin])
    def approve(self, request, pk=None):
        ar = self.get_object()
        if ar.status != 'pending':
            return Response({'detail': 'Déjà traité.'}, status=400)

        with transaction.atomic():
            _apply_approval_request(ar, request)
            ar.status = 'approved'
            ar.decided_by = request.user
            ar.decided_at = timezone.now()
            ar.rejection_reason = None
            ar.save(update_fields=['status', 'decided_by', 'decided_at', 'rejection_reason', 'updated_at'])

        _notify_user(ar.requested_by, 'Action approuvée', f"Votre action ({ar.model}) a été approuvée.")
        return Response(ApprovalRequestSerializer(ar).data)

    @action(detail=True, methods=['post'], url_path='reject', permission_classes=[IsAdminOrSuperAdmin])
    def reject(self, request, pk=None):
        ar = self.get_object()
        if ar.status != 'pending':
            return Response({'detail': 'Déjà traité.'}, status=400)

        reason = request.data.get('reason')
        ar.status = 'rejected'
        ar.decided_by = request.user
        ar.decided_at = timezone.now()
        ar.rejection_reason = (str(reason).strip() if reason is not None else None) or None
        ar.save(update_fields=['status', 'decided_by', 'decided_at', 'rejection_reason', 'updated_at'])

        msg = f"Votre action ({ar.model}) a été refusée.".strip()
        if ar.rejection_reason:
            msg = f"{msg} Motif: {ar.rejection_reason}".strip()
        _notify_user(ar.requested_by, 'Action refusée', msg)
        return Response(ApprovalRequestSerializer(ar).data)


class ChurchBiographyViewSet(viewsets.ModelViewSet):
    """ViewSet pour gérer la biographie de l'église"""
    queryset = ChurchBiography.objects.all()
    serializer_class = ChurchBiographySerializer
    permission_classes = [PublicReadAdminWrite]

    def get_queryset(self):
        qs = super().get_queryset()
        return qs.filter(is_active=True).order_by('-updated_at')

    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)
        self._log_action('create', serializer.instance)

    def perform_update(self, serializer):
        serializer.save()
        self._log_action('update', serializer.instance)

    def destroy(self, request, *args, **kwargs):
        obj = self.get_object()
        obj.is_active = False
        obj.save()
        self._log_action('delete', obj)
        return Response(status=204)

    def _log_action(self, action, obj):
        AuditLogEntry.objects.create(
            actor=getattr(self.request, 'user', None),
            action=action,
            model='ChurchBiography',
            object_id=str(getattr(obj, 'pk', '')),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'pk', '')),
            ip_address=_client_ip(self.request),
            payload={
                'title': getattr(obj, 'title', None),
                'is_active': getattr(obj, 'is_active', None),
            },
        )


class ChurchConsistoryViewSet(viewsets.ModelViewSet):
    """ViewSet pour gérer les informations du consistoire"""
    queryset = ChurchConsistory.objects.all()
    serializer_class = ChurchConsistorySerializer
    permission_classes = [PublicReadAdminWrite]

    def get_queryset(self):
        qs = super().get_queryset()
        return qs.filter(is_active=True).order_by('-updated_at')

    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)
        self._log_action('create', serializer.instance)

    def perform_update(self, serializer):
        serializer.save()
        self._log_action('update', serializer.instance)

    def destroy(self, request, *args, **kwargs):
        obj = self.get_object()
        obj.is_active = False
        obj.save()
        self._log_action('delete', obj)
        return Response(status=204)

    def _log_action(self, action, obj):
        AuditLogEntry.objects.create(
            actor=getattr(self.request, 'user', None),
            action=action,
            model='ChurchConsistory',
            object_id=str(getattr(obj, 'pk', '')),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'pk', '')),
            ip_address=_client_ip(self.request),
            payload={
                'title': getattr(obj, 'title', None),
                'is_active': getattr(obj, 'is_active', None),
            },
        )


class ContactViewSet(viewsets.ModelViewSet):
    """ViewSet pour gérer les messages de contact"""
    queryset = Contact.objects.all()
    serializer_class = ContactSerializer
    permission_classes = [IsAuthenticated]

    def get_permissions(self):
        """Autoriser la création (POST) sans authentification pour le formulaire public"""
        if self.action == 'create':
            return [AllowAny()]
        return super().get_permissions()

    def get_queryset(self):
        qs = super().get_queryset()
        # Filtrer par statut si fourni
        status = self.request.query_params.get('status')
        if status:
            qs = qs.filter(status=status)
        # Filtrer par sujet si fourni
        subject = self.request.query_params.get('subject')
        if subject:
            qs = qs.filter(subject=subject)
        return qs.order_by('-created_at')

    def perform_create(self, serializer):
        """Créer un message avec capture automatique de l'IP et user agent"""
        request = self.request
        ip = _client_ip(request)
        user_agent = request.META.get('HTTP_USER_AGENT', '')
        
        serializer.save(
            ip_address=ip,
            user_agent=user_agent
        )
        self._log_action('create', serializer.instance)

    def perform_update(self, serializer):
        """Marquer comme répondu si le statut change à 'answered'"""
        from django.utils import timezone
        
        if serializer.validated_data.get('status') == 'answered' and serializer.instance.status != 'answered':
            serializer.save(
                answered_by=self.request.user,
                answered_at=timezone.now()
            )
        else:
            serializer.save()
        self._log_action('update', serializer.instance)

    def destroy(self, request, *args, **kwargs):
        """Suppression définitive du message"""
        obj = self.get_object()
        self._log_action('delete', obj)
        obj.delete()
        return Response(status=204)

    @action(detail=True, methods=['post'])
    def mark_as_read(self, request, pk=None):
        """Action personnalisée pour marquer comme lu"""
        obj = self.get_object()
        obj.status = 'read'
        obj.save()
        self._log_action('mark_as_read', obj)
        return Response(ContactSerializer(obj).data)

    @action(detail=True, methods=['post'])
    def mark_as_answered(self, request, pk=None):
        """Action personnalisée pour marquer comme répondu"""
        from django.utils import timezone
        obj = self.get_object()
        obj.status = 'answered'
        obj.answered_by = request.user
        obj.answered_at = timezone.now()
        obj.save()
        self._log_action('mark_as_answered', obj)
        return Response(ContactSerializer(obj).data)

    def _log_action(self, action, obj):
        """Journaliser les actions"""
        try:
            AuditLogEntry.objects.create(
                actor=getattr(self.request, 'user', None),
                action=action,
                model='Contact',
                object_id=str(getattr(obj, 'pk', '')),
                object_repr=f"{getattr(obj, 'name', '')} - {getattr(obj, 'subject', '')}",
                ip_address=_client_ip(self.request),
                payload={
                    'name': getattr(obj, 'name', None),
                    'email': getattr(obj, 'email', None),
                    'subject': getattr(obj, 'subject', None),
                    'status': getattr(obj, 'status', None),
                },
            )
        except Exception:
            pass  # Ignorer les erreurs de logging
